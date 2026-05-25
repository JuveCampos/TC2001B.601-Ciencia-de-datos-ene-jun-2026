"""
Modulo de estilo compartido para las presentaciones de FUNCIONES del curso
TC2001B (Sesiones 16-20). Estilo Tec Monterrey, 4:3, fuente Ubuntu, con
tarjetas de codigo en fuente monoespaciada (Menlo) sobre fondo oscuro.

Builder principal: slide_funcion() -> encabezado con el nombre de la funcion
en fuente de codigo, una oracion de rol, una tarjeta de codigo (fragmento del
ejercicio) y una tarjeta de explicacion.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === Paleta Tec ===
NAVY = RGBColor(0x1E, 0x4C, 0x7D)
LIGHT_BLUE = RGBColor(0xEA, 0xF2, 0xFA)
SOFT_BLUE = RGBColor(0xDD, 0xE5, 0xED)
SOFT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
DARK_GREEN = RGBColor(0x1A, 0x7A, 0x3A)
SOFT_RED = RGBColor(0xFA, 0xDB, 0xD8)
DARK_RED = RGBColor(0xC0, 0x39, 0x2B)
HIGHLIGHT_GREEN = RGBColor(0xF0, 0xFA, 0xF0)
GRAY_TOP = RGBColor(0xF0, 0xF0, 0xF0)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# === Colores del bloque de codigo (estilo editor oscuro) ===
CODE_BG = RGBColor(0x1E, 0x29, 0x36)       # navy muy oscuro
CODE_FG = RGBColor(0xEA, 0xED, 0xF0)       # texto claro
CODE_COMMENT = RGBColor(0x8C, 0xA6, 0x8C)  # comentarios verde-gris
CODE_ACCENT = RGBColor(0x6F, 0xB3, 0xE0)   # azul claro para el caption

UBUNTU = "Ubuntu"
MONO = "Menlo"   # fuente de codigo clasica, nativa de macOS


def nueva_presentacion():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, text, left, top, width, height, *, size=14, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, italic=False, line_spacing=None,
             font=UBUNTU, anchor=None, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_rect(slide, left, top, width, height, fill_color, *, line=None,
             line_w=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_bg(slide, prs, color=WHITE):
    bg = add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, color)
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_header(slide, title, *, width_in=6.0, font=MONO, size=22):
    """Banda navy con el titulo (por defecto en fuente de codigo)."""
    add_rect(slide, 0, Inches(0.3), Inches(width_in), Inches(0.72), NAVY)
    add_text(slide, title, Inches(0.25), Inches(0.3), Inches(width_in - 0.3),
             Inches(0.72), size=size, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=font)


def add_footer(slide, num, titulo_corto):
    add_text(slide, f"{titulo_corto}  -  {num}",
             Inches(6.0), Inches(7.18), Inches(3.8), Inches(0.22),
             size=9, color=GRAY, align=PP_ALIGN.RIGHT, italic=True)


def add_code_card(slide, code_lines, left, top, width, height, *, size=12,
                  caption=None):
    """Tarjeta de codigo: fondo oscuro, fuente monoespaciada, comentarios
    (lineas que empiezan con #) en verde-gris y codigo en claro."""
    if caption:
        add_text(slide, caption, left, top - Inches(0.32), width, Inches(0.28),
                 size=11, color=NAVY, bold=True, font=UBUNTU)
    add_rect(slide, left, top, width, height, CODE_BG)
    box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12),
                                   width - Inches(0.36), height - Inches(0.24))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.08
        run = p.add_run()
        run.text = line if line != "" else " "
        run.font.size = Pt(size)
        run.font.name = MONO
        es_comentario = line.lstrip().startswith("#")
        run.font.color.rgb = CODE_COMMENT if es_comentario else CODE_FG
    return box


def add_card(slide, left, top, width, height, *, fill=LIGHT_BLUE, title=None,
             body=None, title_color=NAVY, body_color=BLACK, title_size=14,
             body_size=12, line_spacing=1.2):
    add_rect(slide, left, top, width, height, fill)
    if title and body:
        add_text(slide, title, left + Inches(0.16), top + Inches(0.1),
                 width - Inches(0.32), Inches(0.4), size=title_size,
                 bold=True, color=title_color)
        add_text(slide, body, left + Inches(0.16), top + Inches(0.52),
                 width - Inches(0.32), height - Inches(0.62),
                 size=body_size, color=body_color, line_spacing=line_spacing)
    elif title:
        add_text(slide, title, left + Inches(0.16), top, width - Inches(0.32),
                 height, size=title_size, bold=True, color=title_color,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=line_spacing)
    elif body:
        add_text(slide, body, left + Inches(0.16), top + Inches(0.1),
                 width - Inches(0.32), height - Inches(0.2), size=body_size,
                 color=body_color, line_spacing=line_spacing)


# ============ SLIDES DE ALTO NIVEL ============

def slide_portada(prs, titulo, subtitulo, autor_info):
    s = _blank(prs)
    add_bg(s, prs)
    add_rect(s, 0, 0, prs.slide_width, Inches(1.1), GRAY_TOP)
    add_text(s, "Tecnológico de Monterrey", Inches(0.4), Inches(0.25),
             Inches(6), Inches(0.6), size=20, color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 0, Inches(1.1), prs.slide_width, Inches(3.6), SOFT_BLUE)
    add_text(s, titulo, Inches(0.6), Inches(1.7), Inches(8.8), Inches(2.0),
             size=40, color=NAVY, bold=True, line_spacing=1.05)
    add_text(s, subtitulo, Inches(0.6), Inches(3.95), Inches(8.8), Inches(0.7),
             size=18, color=NAVY, italic=True)
    add_text(s, autor_info, Inches(0.6), Inches(5.4), Inches(8.8), Inches(1.0),
             size=14, color=BLACK, line_spacing=1.3)
    add_rect(s, Inches(0.6), Inches(6.7), Inches(2), Inches(0.05), NAVY)
    return s


def slide_intro(prs, titulo, parrafo, pasos, footer_num, footer_txt):
    """Slide de contexto: parrafo + tarjeta con el flujo/lista de pasos."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    add_text(s, parrafo, Inches(0.4), Inches(1.25), Inches(9.2), Inches(1.4),
             size=15, color=BLACK, line_spacing=1.3)
    add_card(s, Inches(0.3), Inches(2.8), Inches(9.4), Inches(3.7),
             fill=LIGHT_BLUE, title="Funciones que veremos, en orden de uso",
             body=pasos, title_size=16, body_size=13, line_spacing=1.45)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_funcion(prs, nombre, rol, codigo, explica, footer_num, footer_txt,
                  *, code_caption="Uso en el ejercicio", code_size=12,
                  explica_titulo="Qué hace"):
    """Slide de una funcion: nombre (fuente codigo) + rol + codigo + explica."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, nombre, width_in=6.0, font=MONO, size=22)
    add_text(s, rol, Inches(0.4), Inches(1.18), Inches(9.2), Inches(0.62),
             size=15, color=NAVY, italic=True, line_spacing=1.2)
    n = len(codigo)
    code_h = max(1.7, min(3.0, 0.32 + n * 0.245))
    add_code_card(s, codigo, Inches(0.3), Inches(2.15), Inches(9.4),
                  Inches(code_h), size=code_size, caption=code_caption)
    exp_top = 2.15 + code_h + 0.55
    exp_h = 7.0 - exp_top
    add_card(s, Inches(0.3), Inches(exp_top), Inches(9.4), Inches(exp_h),
             fill=LIGHT_BLUE, title=explica_titulo, body=explica,
             title_size=15, body_size=13, line_spacing=1.32)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_conclusiones(prs, titulo, items, footer_num, footer_txt):
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    add_card(s, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.4),
             fill=HIGHLIGHT_GREEN, title="Ideas clave para llevarte",
             body=items, title_color=DARK_GREEN, title_size=18, body_size=15,
             line_spacing=1.5)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_cierre(prs):
    s = _blank(prs)
    add_bg(s, prs)
    add_rect(s, 0, 0, prs.slide_width, Inches(1.1), GRAY_TOP)
    add_text(s, "Tecnológico de Monterrey", Inches(0.4), Inches(0.25),
             Inches(6), Inches(0.6), size=20, color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 0, Inches(1.1), prs.slide_width, Inches(3.6), SOFT_BLUE)
    add_text(s, "Gracias", Inches(0.6), Inches(1.8), Inches(8.8), Inches(2),
             size=80, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.0), Inches(5.5), Inches(4.0), Inches(1.0), NAVY)
    add_text(s, "¿Preguntas?", Inches(3.0), Inches(5.5), Inches(4.0),
             Inches(1.0), size=28, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s
