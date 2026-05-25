"""Plantillas de diapositiva para las presentaciones TEORICAS del curso
TC2001B (Sesiones 16-20), pensadas para audiencia de preparatoria: muchas
figuras, esquemas numerados, analogias cotidianas y comparaciones.

Reutiliza la paleta y los ayudantes base de estilo_tec.py y agrega plantillas
con imagenes (las figuras se generan con figteoria.py). Estilo Tec 4:3, Ubuntu.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from estilo_tec import (
    nueva_presentacion, _blank, add_text, add_rect, add_bg, add_header,
    add_footer, add_card, slide_portada, slide_conclusiones, slide_cierre,
    NAVY, LIGHT_BLUE, SOFT_BLUE, SOFT_GREEN, DARK_GREEN, SOFT_RED, DARK_RED,
    HIGHLIGHT_GREEN, GRAY_TOP, GRAY, WHITE, BLACK, UBUNTU,
)
from pptx.dml.color import RGBColor

# Tarjetas planas del curso (acentos pedagogicos)
TEAL = RGBColor(0x2D, 0x7A, 0x7A)
OCHRE = RGBColor(0xC9, 0x7B, 0x2A)
PLUM = RGBColor(0x7D, 0x3C, 0x66)
OLIVE = RGBColor(0x5A, 0x7D, 0x2A)
SLATE = RGBColor(0x4A, 0x55, 0x68)
SOFT_OCHRE = RGBColor(0xF7, 0xEC, 0xDD)
SOFT_TEAL = RGBColor(0xDD, 0xEC, 0xEC)


# ---------------------------------------------------------------------------
# Ayudantes
# ---------------------------------------------------------------------------
def add_picture_fit(slide, path, left, top, width, height):
    """Inserta una imagen ajustada a la caja (left, top, width, height)
    preservando su proporcion y centrandola."""
    pic = slide.shapes.add_picture(path, left, top, width=width)
    if pic.height > height:
        slide.shapes._spTree.remove(pic._element)
        pic = slide.shapes.add_picture(path, left, top, height=height)
    pic.left = int(left + (width - pic.width) / 2)
    pic.top = int(top + (height - pic.height) / 2)
    return pic


def _badge(slide, n, left, top, d=Inches(0.42), color=NAVY):
    """Circulo de color con un numero blanco (tarjeta numerada del curso)."""
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    o.fill.solid()
    o.fill.fore_color.rgb = color
    o.line.fill.background()
    o.shadow.inherit = False
    tf = o.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(n)
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = UBUNTU
    return o


def add_fuente(slide, texto):
    """Cita de fuente al pie (9pt gris)."""
    add_text(slide, texto, Inches(0.4), Inches(7.0), Inches(9.2), Inches(0.24),
             size=9, color=GRAY, italic=True)


# ---------------------------------------------------------------------------
# Plantillas de seccion / portada interna
# ---------------------------------------------------------------------------
def slide_seccion(prs, numero, titulo, subtitulo=""):
    """Separador de seccion: fondo navy, numero grande y titulo en blanco."""
    s = _blank(prs)
    add_bg(s, prs, NAVY)
    add_text(s, str(numero), Inches(0.6), Inches(1.4), Inches(3.2), Inches(3.0),
             size=160, bold=True, color=RGBColor(0x33, 0x63, 0x95),
             align=PP_ALIGN.LEFT)
    add_rect(s, Inches(3.7), Inches(2.4), Inches(0.08), Inches(2.2), WHITE)
    add_text(s, titulo, Inches(4.0), Inches(2.5), Inches(5.6), Inches(1.8),
             size=34, bold=True, color=WHITE, line_spacing=1.05,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitulo:
        add_text(s, subtitulo, Inches(4.0), Inches(4.45), Inches(5.6),
                 Inches(0.9), size=15, color=RGBColor(0xC7, 0xD7, 0xE6),
                 italic=True, line_spacing=1.2)
    return s


def slide_agenda(prs, titulo, items, footer_num, footer_txt):
    """Mapa del tema: lista de secciones con badges numerados en dos columnas."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    n = len(items)
    mitad = (n + 1) // 2
    cols = [(Inches(0.5), items[:mitad], 0), (Inches(5.15), items[mitad:], mitad)]
    for x, grupo, offset in cols:
        for i, it in enumerate(grupo):
            top = Inches(1.45 + i * 0.95)
            _badge(s, offset + i + 1, x, top, d=Inches(0.5))
            add_text(s, it, x + Inches(0.68), top - Inches(0.02), Inches(3.7),
                     Inches(0.9), size=14, color=BLACK, line_spacing=1.1,
                     anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, footer_num, footer_txt)
    return s


# ---------------------------------------------------------------------------
# Plantillas de contenido
# ---------------------------------------------------------------------------
def slide_concepto(prs, titulo, definicion, footer_num, footer_txt, *,
                   en_cristiano=None, etiqueta="Definición"):
    """Concepto clave: tarjeta grande con la definicion y, opcional, una
    reformulacion en lenguaje llano (caja verde 'En cristiano')."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    h_def = Inches(3.0) if en_cristiano else Inches(4.9)
    add_card(s, Inches(0.3), Inches(1.35), Inches(9.4), h_def, fill=LIGHT_BLUE,
             title=etiqueta, body=definicion, title_size=18, body_size=18,
             line_spacing=1.4)
    if en_cristiano:
        add_card(s, Inches(0.3), Inches(4.6), Inches(9.4), Inches(2.2),
                 fill=HIGHLIGHT_GREEN, title="En palabras simples",
                 body=en_cristiano, title_color=DARK_GREEN, title_size=16,
                 body_size=16, line_spacing=1.4)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_figura(prs, titulo, img, footer_num, footer_txt, *, caption=None,
                 fuente=None, nota=None):
    """Figura grande centrada + caption opcional + nota opcional."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    img_h = Inches(4.4) if nota else Inches(5.0)
    add_picture_fit(s, img, Inches(0.5), Inches(1.35), Inches(9.0), img_h)
    y = 1.35 + (4.4 if nota else 5.0) + 0.05
    if caption:
        add_text(s, caption, Inches(0.5), Inches(y), Inches(9.0), Inches(0.4),
                 size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    if nota:
        add_card(s, Inches(0.3), Inches(6.0), Inches(9.4), Inches(0.95),
                 fill=HIGHLIGHT_GREEN, body=nota, body_size=14,
                 body_color=BLACK, line_spacing=1.25)
    if fuente:
        add_fuente(s, fuente)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_texto_figura(prs, titulo, intro, puntos, img, footer_num, footer_txt,
                       *, img_side="right", caption=None, fuente=None,
                       puntos_size=14):
    """Dos paneles: texto (intro + puntos en tarjeta) y figura al lado."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    if intro:
        add_text(s, intro, Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.7),
                 size=15, color=NAVY, italic=True, line_spacing=1.2)
    top = Inches(2.05)
    txt_w, img_w = Inches(4.5), Inches(4.6)
    if img_side == "right":
        tx, ix = Inches(0.3), Inches(5.1)
    else:
        tx, ix = Inches(5.2), Inches(0.3)
    add_card(s, tx, top, txt_w, Inches(4.5), fill=LIGHT_BLUE, body=puntos,
             body_size=puntos_size, line_spacing=1.4)
    add_picture_fit(s, img, ix, top, img_w, Inches(4.5))
    if caption:
        add_text(s, caption, ix, Inches(6.6), img_w, Inches(0.34), size=11,
                 color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    if fuente:
        add_fuente(s, fuente)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_pasos(prs, titulo, pasos, footer_num, footer_txt, *, intro=None,
                img=None):
    """Esquema numerado: lista de (titulo_paso, descripcion) con badges.
    Si hay img, los pasos van a la izquierda y la figura a la derecha."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    y0 = 1.25
    if intro:
        add_text(s, intro, Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.6),
                 size=14, color=NAVY, italic=True, line_spacing=1.2)
        y0 = 1.95
    col_w = Inches(4.5) if img else Inches(9.0)
    n = len(pasos)
    paso_h = min(1.15, (6.7 - y0) / max(n, 1))
    for i, (pt, pd) in enumerate(pasos):
        top = Inches(y0 + i * paso_h)
        _badge(s, i + 1, Inches(0.45), top, d=Inches(0.5))
        add_text(s, pt, Inches(1.15), top - Inches(0.04), col_w - Inches(0.9),
                 Inches(0.4), size=14, bold=True, color=NAVY)
        if pd:
            add_text(s, pd, Inches(1.15), top + Inches(0.32),
                     col_w - Inches(0.9), Inches(paso_h - 0.34), size=12,
                     color=BLACK, line_spacing=1.15)
    if img:
        add_picture_fit(s, img, Inches(5.1), Inches(y0), Inches(4.6),
                        Inches(6.6 - y0))
    add_footer(s, footer_num, footer_txt)
    return s


def slide_comparacion(prs, titulo, col_izq, col_der, footer_num, footer_txt,
                      *, intro=None):
    """Dos columnas con figura + cuerpo. Cada col es dict con claves:
    titulo, img (opcional), body, color (encabezado de la tarjeta), fill."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    y0 = 1.3
    if intro:
        add_text(s, intro, Inches(0.4), Inches(1.18), Inches(9.2), Inches(0.55),
                 size=14, color=NAVY, italic=True, line_spacing=1.15)
        y0 = 1.85
    for col, x in ((col_izq, Inches(0.3)), (col_der, Inches(5.05))):
        w = Inches(4.65)
        fill = col.get("fill") or LIGHT_BLUE
        hdr_color = col.get("color") or NAVY
        add_rect(s, x, Inches(y0), w, Inches(0.55), hdr_color)
        add_text(s, col["titulo"], x, Inches(y0), w, Inches(0.55), size=15,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        cy = y0 + 0.55
        if col.get("img"):
            add_picture_fit(s, col["img"], x, Inches(cy + 0.08), w,
                            Inches(2.7))
            cy += 2.85
        add_rect(s, x, Inches(cy), w, Inches(6.55 - cy), fill)
        add_text(s, col["body"], x + Inches(0.16), Inches(cy + 0.12),
                 w - Inches(0.32), Inches(6.4 - cy), size=13, color=BLACK,
                 line_spacing=1.3)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_analogia(prs, titulo, texto, footer_num, footer_txt, *, img=None,
                   color=OCHRE, fill=SOFT_OCHRE, etiqueta="Piénsalo así"):
    """Analogia cotidiana: tarjeta de acento (ocre por defecto) con la idea
    en lenguaje llano; opcionalmente una figura al lado."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    if img:
        add_rect(s, Inches(0.3), Inches(1.4), Inches(0.12), Inches(5.0), color)
        add_card(s, Inches(0.5), Inches(1.4), Inches(4.6), Inches(5.0),
                 fill=fill, title=etiqueta, body=texto, title_color=color,
                 title_size=18, body_size=16, line_spacing=1.45)
        add_picture_fit(s, img, Inches(5.3), Inches(1.4), Inches(4.4),
                        Inches(5.0))
    else:
        add_rect(s, Inches(0.3), Inches(1.55), Inches(0.14), Inches(4.6), color)
        add_card(s, Inches(0.5), Inches(1.55), Inches(9.2), Inches(4.6),
                 fill=fill, title=etiqueta, body=texto, title_color=color,
                 title_size=20, body_size=18, line_spacing=1.5)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_alerta(prs, titulo, mensaje, footer_num, footer_txt, *, img=None,
                 etiqueta="Error común"):
    """Diapositiva critica: banner rojo + tarjeta roja (errores, sesgos)."""
    s = _blank(prs)
    add_bg(s, prs)
    add_rect(s, 0, Inches(0.3), Inches(9.4), Inches(0.72), DARK_RED)
    add_text(s, titulo, Inches(0.25), Inches(0.3), Inches(9.1), Inches(0.72),
             size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if img:
        add_card(s, Inches(0.3), Inches(1.4), Inches(4.6), Inches(5.0),
                 fill=SOFT_RED, title=etiqueta, body=mensaje,
                 title_color=DARK_RED, title_size=18, body_size=15,
                 line_spacing=1.4)
        add_picture_fit(s, img, Inches(5.3), Inches(1.4), Inches(4.4),
                        Inches(5.0))
    else:
        add_card(s, Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.7),
                 fill=SOFT_RED, title=etiqueta, body=mensaje,
                 title_color=DARK_RED, title_size=18, body_size=17,
                 line_spacing=1.5)
    add_footer(s, footer_num, footer_txt)
    return s


def slide_tabla(prs, titulo, encabezados, filas, footer_num, footer_txt, *,
                intro=None, col_widths=None):
    """Tabla simple estilo Tec (encabezado navy, filas alternadas)."""
    s = _blank(prs)
    add_bg(s, prs)
    add_header(s, titulo, width_in=9.4, font=UBUNTU, size=22)
    top = 1.4
    if intro:
        add_text(s, intro, Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.6),
                 size=14, color=NAVY, italic=True, line_spacing=1.2)
        top = 2.0
    ncol = len(encabezados)
    nrow = len(filas) + 1
    tbl_w, tbl_h = Inches(9.4), Inches(min(5.0, 0.55 * nrow + 0.2))
    gf = s.shapes.add_table(nrow, ncol, Inches(0.3), Inches(top), tbl_w,
                            tbl_h).table
    if col_widths:
        for j, w in enumerate(col_widths):
            gf.columns[j].width = Inches(w)
    for j, h in enumerate(encabezados):
        c = gf.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        _set_cell(c, h, WHITE, bold=True, size=13)
    for i, fila in enumerate(filas, start=1):
        for j, val in enumerate(fila):
            c = gf.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT_BLUE if i % 2 else WHITE
            _set_cell(c, val, BLACK, bold=False, size=12)
    add_footer(s, footer_num, footer_txt)
    return s


def _set_cell(cell, text, color, *, bold, size):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = UBUNTU
