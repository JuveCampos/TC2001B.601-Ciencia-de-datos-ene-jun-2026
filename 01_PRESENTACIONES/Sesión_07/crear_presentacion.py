#!/usr/bin/env python3
"""
Genera la presentación PowerPoint – Fichas Sesión 07 (modo CLARO).
Estilo compatible con Clase07.key: fondo blanco, texto oscuro, acentos en azul.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Paleta modo claro ─────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)  # fondo principal
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)  # texto principal (casi negro)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x66)  # texto secundario
LIGHT_GRAY  = RGBColor(0xE0, 0xE0, 0xE8)  # separadores suaves
MID_GRAY    = RGBColor(0xCC, 0xCC, 0xD5)  # divisores

# Acentos (saturados para legibilidad en fondo claro)
BLUE_TEC    = RGBColor(0x00, 0x3D, 0xA5)  # azul Tec de Monterrey
ACCENT_BLUE = RGBColor(0x1A, 0x6B, 0xD4)  # azul funciones
ACCENT_GREEN= RGBColor(0x1B, 0x8A, 0x4A)  # verde
ACCENT_PINK = RGBColor(0xC0, 0x39, 0x7A)  # rosa fuerte
ACCENT_PEACH= RGBColor(0xD0, 0x6B, 0x21)  # naranja/durazno
ACCENT_YELLOW=RGBColor(0xA8, 0x83, 0x07)  # amarillo oscuro (legible)
ACCENT_MAUVE= RGBColor(0x7C, 0x3A, 0xED)  # morado
ACCENT_RED  = RGBColor(0xDC, 0x26, 0x26)  # rojo
ACCENT_TEAL = RGBColor(0x0D, 0x7C, 0x7C)  # teal oscuro

CODE_BG     = RGBColor(0xF3, 0xF4, 0xF8)  # gris muy claro para code blocks
CODE_TEXT   = RGBColor(0x1E, 0x29, 0x3B)  # texto del código
CODE_BORDER = RGBColor(0xD0, 0xD5, 0xDD)  # borde del code block
BADGE_TEXT  = RGBColor(0xFF, 0xFF, 0xFF)  # texto dentro de badges

# Footer
FOOTER_BG   = RGBColor(0x00, 0x3D, 0xA5)  # barra inferior azul Tec

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

# ─── Helpers ───────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_footer_bar(slide):
    """Barra inferior azul tipo Tec de Monterrey."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(8.6), Inches(16), Inches(0.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOOTER_BG
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]
    p.text = "TC2001B.601 · Ciencia de datos · Sesión 07"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.name = "Ubuntu"
    p.alignment = PP_ALIGN.LEFT

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Ubuntu"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_paragraph(tf, text, font_size=16, color=DARK_TEXT, bold=False,
                  font_name="Ubuntu", alignment=PP_ALIGN.LEFT, space_after=Pt(6)):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p

def add_code_block(slide, left, top, width, height, code_text, font_size=13):
    """Rectángulo gris claro con código mono."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = CODE_BORDER
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.02

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_right = Inches(0.3)
    tf.margin_bottom = Inches(0.2)

    lines = code_text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = CODE_TEXT
        p.font.name = "Ubuntu Mono"
        p.space_after = Pt(2)
    return shape

def add_badge(slide, left, top, width, height, text, bg_color, text_color=BADGE_TEXT, font_size=12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.15
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Ubuntu"
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    return shape

def add_function_slide(func_name, package_name, description, syntax, example_code,
                       notes, accent_color, tip=None):
    """Plantilla estándar para fichas de funciones – modo claro."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    # Badge del paquete
    add_badge(slide, 0.5, 0.4, 2.0, 0.45, package_name, accent_color)

    # Nombre de la función (título)
    add_textbox(slide, 0.5, 1.0, 10, 0.8, func_name,
                font_size=40, color=accent_color, bold=True, font_name="Ubuntu Mono")

    # Descripción
    tf_desc = add_rich_textbox(slide, 0.5, 1.9, 7.0, 1.5)
    add_paragraph(tf_desc, description, font_size=18, color=GRAY_TEXT, font_name="Ubuntu")

    # Sintaxis
    add_textbox(slide, 0.5, 3.4, 3, 0.4, "Sintaxis:", font_size=14,
                color=accent_color, bold=True, font_name="Ubuntu")
    add_code_block(slide, 0.5, 3.8, 7.0, 1.0, syntax, font_size=14)

    # Ejemplo
    add_textbox(slide, 0.5, 5.0, 3, 0.4, "Ejemplo:", font_size=14,
                color=accent_color, bold=True, font_name="Ubuntu")
    add_code_block(slide, 0.5, 5.4, 7.0, 2.5, example_code, font_size=13)

    # Panel derecho con notas
    add_textbox(slide, 8.5, 1.0, 7, 0.5, "Notas clave",
                font_size=22, color=DARK_TEXT, bold=True, font_name="Ubuntu")

    # Línea separadora
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.03)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = accent_color
    sep.line.fill.background()

    tf_notes = add_rich_textbox(slide, 8.5, 1.7, 6.8, 5.5)
    for note in notes:
        add_paragraph(tf_notes, f"  {note}", font_size=16, color=GRAY_TEXT, font_name="Ubuntu",
                      space_after=Pt(10))

    if tip:
        # Caja tip
        tip_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.5), Inches(7.2), Inches(6.8), Inches(1.1)
        )
        tip_shape.fill.solid()
        tip_shape.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)  # amarillo pastel
        tip_shape.line.color.rgb = RGBColor(0xF5, 0xD0, 0x50)
        tip_shape.line.width = Pt(1)
        tip_shape.adjustments[0] = 0.05
        tf_tip = tip_shape.text_frame
        tf_tip.word_wrap = True
        tf_tip.margin_left = Inches(0.2)
        tf_tip.margin_top = Inches(0.1)
        p_tip_title = tf_tip.paragraphs[0]
        p_tip_title.text = "Tip"
        p_tip_title.font.size = Pt(13)
        p_tip_title.font.bold = True
        p_tip_title.font.color.rgb = ACCENT_PEACH
        p_tip_title.font.name = "Ubuntu"
        p_tip = tf_tip.add_paragraph()
        p_tip.text = tip
        p_tip.font.size = Pt(13)
        p_tip.font.color.rgb = GRAY_TEXT
        p_tip.font.name = "Ubuntu"

    return slide


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 1: geom_col() - Gráfica de barras
# ═══════════════════════════════════════════════════════════════════════════
def slide_geom_col():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 2.0, 0.45, "ggplot2", ACCENT_BLUE)
    add_textbox(slide, 0.5, 1.0, 10, 0.8, "geom_col()",
                font_size=40, color=ACCENT_BLUE, bold=True, font_name="Ubuntu Mono")

    tf = add_rich_textbox(slide, 0.5, 1.9, 7.0, 1.2)
    add_paragraph(tf, "Crea gráficas de barras donde la altura de cada barra "
                  "representa un valor ya calculado en los datos.",
                  font_size=18, color=GRAY_TEXT)
    add_paragraph(tf, "A diferencia de geom_bar(), NO cuenta frecuencias.",
                  font_size=16, color=ACCENT_RED, bold=True)

    add_textbox(slide, 0.5, 3.3, 3, 0.4, "Sintaxis:", font_size=14,
                color=ACCENT_BLUE, bold=True)
    add_code_block(slide, 0.5, 3.7, 7.0, 0.9,
        "ggplot(datos, aes(x = categoria, y = valor)) +\n"
        "  geom_col()", font_size=14)

    add_textbox(slide, 0.5, 4.8, 3, 0.4, "Ejemplo:", font_size=14,
                color=ACCENT_BLUE, bold=True)
    add_code_block(slide, 0.5, 5.2, 7.0, 3.0,
        "library(tidyverse)\n\n"
        "# Datos de ventas por mes\n"
        "ventas <- tibble(\n"
        '  mes = c("Ene", "Feb", "Mar", "Abr"),\n'
        "  total = c(15000, 22000, 18000, 25000)\n"
        ")\n\n"
        "ggplot(ventas, aes(x = mes, y = total)) +\n"
        '  geom_col(fill = "#1A6BD4") +\n'
        '  labs(title = "Ventas mensuales")',
        font_size=12)

    # Panel derecho
    add_textbox(slide, 8.5, 1.0, 7, 0.5, "Parámetros comunes",
                font_size=22, color=DARK_TEXT, bold=True)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.03))
    sep.fill.solid(); sep.fill.fore_color.rgb = ACCENT_BLUE; sep.line.fill.background()

    params = [
        ("fill", "Color de relleno de las barras"),
        ("color", "Color del borde de las barras"),
        ("width", "Ancho de las barras (0 a 1)"),
        ("alpha", "Transparencia (0 = invisible, 1 = sólido)"),
        ("position", '"dodge" para barras lado a lado,\n"stack" para apiladas'),
    ]
    y_pos = 1.7
    for param, desc in params:
        add_textbox(slide, 8.5, y_pos, 3, 0.35, param,
                    font_size=16, color=ACCENT_BLUE, bold=True, font_name="Ubuntu Mono")
        add_textbox(slide, 10.5, y_pos, 5, 0.5, desc,
                    font_size=14, color=GRAY_TEXT)
        y_pos += 0.55

    # geom_col vs geom_bar
    add_textbox(slide, 8.5, 4.8, 7, 0.4, "geom_col() vs geom_bar()",
                font_size=18, color=DARK_TEXT, bold=True)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(5.2), Inches(6.5), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = ACCENT_BLUE; sep2.line.fill.background()

    tf2 = add_rich_textbox(slide, 8.5, 5.4, 6.8, 2.5)
    add_paragraph(tf2, "geom_col():", font_size=15, color=ACCENT_BLUE, bold=True)
    add_paragraph(tf2, "  Usa los valores tal como están en los datos.", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf2, "  Requiere aes(x, y).", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf2, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf2, "geom_bar():", font_size=15, color=ACCENT_PEACH, bold=True)
    add_paragraph(tf2, "  Cuenta frecuencias automáticamente.", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf2, "  Solo necesita aes(x).", font_size=14, color=GRAY_TEXT)

slide_geom_col()


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 2: geom_point()
# ═══════════════════════════════════════════════════════════════════════════
def slide_geom_point():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 2.0, 0.45, "ggplot2", ACCENT_GREEN)
    add_textbox(slide, 0.5, 1.0, 10, 0.8, "geom_point()",
                font_size=40, color=ACCENT_GREEN, bold=True, font_name="Ubuntu Mono")

    tf = add_rich_textbox(slide, 0.5, 1.9, 7.0, 1.2)
    add_paragraph(tf, "Crea gráficas de dispersión (scatterplots). "
                  "Ideal para explorar relaciones entre dos variables numéricas.",
                  font_size=18, color=GRAY_TEXT)

    add_textbox(slide, 0.5, 3.3, 3, 0.4, "Sintaxis:", font_size=14,
                color=ACCENT_GREEN, bold=True)
    add_code_block(slide, 0.5, 3.7, 7.0, 0.9,
        "ggplot(datos, aes(x = var_x, y = var_y)) +\n"
        "  geom_point()", font_size=14)

    add_textbox(slide, 0.5, 4.8, 3, 0.4, "Ejemplo:", font_size=14,
                color=ACCENT_GREEN, bold=True)
    add_code_block(slide, 0.5, 5.2, 7.0, 3.0,
        "library(tidyverse)\n\n"
        "# Relación entre cilindrada y rendimiento\n"
        "ggplot(mpg, aes(x = displ, y = hwy)) +\n"
        "  geom_point(\n"
        '    color = "#1B8A4A",\n'
        "    size = 3,\n"
        "    alpha = 0.7\n"
        "  ) +\n"
        "  labs(\n"
        '    title = "Cilindrada vs Rendimiento",\n'
        '    x = "Cilindrada (litros)",\n'
        '    y = "Millas por galón (autopista)"\n'
        "  )",
        font_size=12)

    # Panel derecho
    add_textbox(slide, 8.5, 1.0, 7, 0.5, "Parámetros comunes",
                font_size=22, color=DARK_TEXT, bold=True)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.03))
    sep.fill.solid(); sep.fill.fore_color.rgb = ACCENT_GREEN; sep.line.fill.background()

    params = [
        ("color", "Color de los puntos"),
        ("size", "Tamaño de los puntos (default ~1.5)"),
        ("shape", "Forma: 16=círculo, 17=triángulo, 15=cuadrado..."),
        ("alpha", "Transparencia (útil si hay sobreposición)"),
        ("stroke", "Grosor del borde (para shapes con borde)"),
    ]
    y_pos = 1.7
    for param, desc in params:
        add_textbox(slide, 8.5, y_pos, 3, 0.35, param,
                    font_size=16, color=ACCENT_GREEN, bold=True, font_name="Ubuntu Mono")
        add_textbox(slide, 10.5, y_pos, 5, 0.5, desc,
                    font_size=14, color=GRAY_TEXT)
        y_pos += 0.55

    # Mapeos estéticos
    add_textbox(slide, 8.5, 4.8, 7, 0.4, "Mapeos estéticos dentro de aes()",
                font_size=18, color=DARK_TEXT, bold=True)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(5.2), Inches(6.5), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = ACCENT_GREEN; sep2.line.fill.background()

    add_code_block(slide, 8.5, 5.4, 6.8, 2.5,
        "# Color por categoría\n"
        "aes(color = clase)\n\n"
        "# Tamaño por variable numérica\n"
        "aes(size = poblacion)\n\n"
        "# Forma por grupo\n"
        "aes(shape = tipo)",
        font_size=13)

slide_geom_point()


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 3: geom_line()
# ═══════════════════════════════════════════════════════════════════════════
def slide_geom_line():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 2.0, 0.45, "ggplot2", ACCENT_PINK)
    add_textbox(slide, 0.5, 1.0, 10, 0.8, "geom_line()",
                font_size=40, color=ACCENT_PINK, bold=True, font_name="Ubuntu Mono")

    tf = add_rich_textbox(slide, 0.5, 1.9, 7.0, 1.2)
    add_paragraph(tf, "Conecta observaciones con líneas, ordenadas por el eje X. "
                  "Ideal para series de tiempo y datos con progresión.",
                  font_size=18, color=GRAY_TEXT)

    add_textbox(slide, 0.5, 3.3, 3, 0.4, "Sintaxis:", font_size=14,
                color=ACCENT_PINK, bold=True)
    add_code_block(slide, 0.5, 3.7, 7.0, 0.9,
        "ggplot(datos, aes(x = tiempo, y = valor)) +\n"
        "  geom_line()", font_size=14)

    add_textbox(slide, 0.5, 4.8, 3, 0.4, "Ejemplo:", font_size=14,
                color=ACCENT_PINK, bold=True)
    add_code_block(slide, 0.5, 5.2, 7.0, 3.0,
        "library(tidyverse)\n\n"
        "# Serie de tiempo de temperatura\n"
        "temps <- tibble(\n"
        '  fecha = seq(as.Date("2025-01-01"),\n'
        '              as.Date("2025-12-01"), by = "month"),\n'
        "  temp = c(12,14,18,22,28,32,34,33,29,23,17,13)\n"
        ")\n\n"
        "ggplot(temps, aes(x = fecha, y = temp)) +\n"
        '  geom_line(color = "#C0397A", linewidth = 1.2) +\n'
        '  geom_point(color = "#C0397A", size = 3) +\n'
        '  labs(title = "Temperatura mensual 2025")',
        font_size=11)

    # Panel derecho
    add_textbox(slide, 8.5, 1.0, 7, 0.5, "Parámetros comunes",
                font_size=22, color=DARK_TEXT, bold=True)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.03))
    sep.fill.solid(); sep.fill.fore_color.rgb = ACCENT_PINK; sep.line.fill.background()

    params = [
        ("color", "Color de la línea"),
        ("linewidth", "Grosor de la línea (antes era 'size')"),
        ("linetype", '"solid", "dashed", "dotted", "dotdash"...'),
        ("alpha", "Transparencia de la línea"),
        ("group", "Agrupa líneas (dentro de aes())"),
    ]
    y_pos = 1.7
    for param, desc in params:
        add_textbox(slide, 8.5, y_pos, 3, 0.35, param,
                    font_size=16, color=ACCENT_PINK, bold=True, font_name="Ubuntu Mono")
        add_textbox(slide, 10.5, y_pos, 5, 0.5, desc,
                    font_size=14, color=GRAY_TEXT)
        y_pos += 0.55

    # Tip
    add_textbox(slide, 8.5, 4.8, 7, 0.4, "Tip: Combina geom_line + geom_point",
                font_size=18, color=DARK_TEXT, bold=True)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(5.2), Inches(6.5), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = ACCENT_PINK; sep2.line.fill.background()

    tf2 = add_rich_textbox(slide, 8.5, 5.4, 6.8, 2.5)
    add_paragraph(tf2, "Agregar puntos sobre la línea ayuda a identificar "
                  "las observaciones exactas:", font_size=15, color=GRAY_TEXT)
    add_paragraph(tf2, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf2, "  geom_line() + geom_point()", font_size=16,
                  color=ACCENT_PINK, bold=True, font_name="Ubuntu Mono")
    add_paragraph(tf2, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf2, "Para múltiples líneas, usa:", font_size=15, color=GRAY_TEXT)
    add_paragraph(tf2, "  aes(color = grupo)", font_size=16,
                  color=ACCENT_PINK, bold=True, font_name="Ubuntu Mono")

slide_geom_line()


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 4: geom_sf()
# ═══════════════════════════════════════════════════════════════════════════
def slide_geom_sf():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 2.2, 0.45, "ggplot2 + sf", ACCENT_TEAL)
    add_textbox(slide, 0.5, 1.0, 10, 0.8, "geom_sf()",
                font_size=40, color=ACCENT_TEAL, bold=True, font_name="Ubuntu Mono")

    tf = add_rich_textbox(slide, 0.5, 1.9, 7.0, 1.2)
    add_paragraph(tf, "Dibuja geometrías espaciales (mapas) a partir de "
                  "objetos sf (Simple Features). Polígonos, puntos, líneas.",
                  font_size=18, color=GRAY_TEXT)

    add_textbox(slide, 0.5, 3.3, 3, 0.4, "Sintaxis:", font_size=14,
                color=ACCENT_TEAL, bold=True)
    add_code_block(slide, 0.5, 3.7, 7.0, 0.9,
        "ggplot(datos_sf) +\n"
        "  geom_sf()", font_size=14)

    add_textbox(slide, 0.5, 4.8, 3, 0.4, "Ejemplo:", font_size=14,
                color=ACCENT_TEAL, bold=True)
    add_code_block(slide, 0.5, 5.2, 7.0, 3.0,
        "library(tidyverse)\n"
        "library(sf)\n\n"
        "# Leer un shapefile de estados de México\n"
        'mapa <- st_read("estados.shp")\n\n'
        "# Mapa coloreado por población\n"
        "ggplot(mapa) +\n"
        '  geom_sf(aes(fill = poblacion),\n'
        '          color = "white", linewidth = 0.2) +\n'
        "  scale_fill_viridis_c() +\n"
        "  theme_void() +\n"
        '  labs(title = "Población por estado")',
        font_size=11)

    # Panel derecho
    add_textbox(slide, 8.5, 1.0, 7, 0.5, "Notas clave",
                font_size=22, color=DARK_TEXT, bold=True)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.03))
    sep.fill.solid(); sep.fill.fore_color.rgb = ACCENT_TEAL; sep.line.fill.background()

    tf2 = add_rich_textbox(slide, 8.5, 1.7, 6.8, 3.0)
    notes = [
        "Requiere el paquete sf instalado",
        "Los datos deben ser un objeto sf con columna geometry",
        "ggplot detecta automáticamente las coordenadas",
        "fill colorea el interior de los polígonos",
        "color controla el borde de los polígonos",
        "Usa theme_void() para un mapa más limpio",
        "scale_fill_viridis_c() es ideal para mapas",
    ]
    for note in notes:
        add_paragraph(tf2, f"  {note}", font_size=15, color=GRAY_TEXT, space_after=Pt(8))

    # Formatos de datos
    add_textbox(slide, 8.5, 5.0, 7, 0.4, "Fuentes de datos geográficos",
                font_size=18, color=DARK_TEXT, bold=True)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(5.4), Inches(6.5), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = ACCENT_TEAL; sep2.line.fill.background()

    tf3 = add_rich_textbox(slide, 8.5, 5.6, 6.8, 2.5)
    add_paragraph(tf3, ".shp (Shapefile) — El más común", font_size=14, color=ACCENT_TEAL, bold=True)
    add_paragraph(tf3, ".geojson — Formato web, texto plano", font_size=14, color=ACCENT_TEAL, bold=True)
    add_paragraph(tf3, ".gpkg (GeoPackage) — Moderno, un solo archivo", font_size=14, color=ACCENT_TEAL, bold=True)
    add_paragraph(tf3, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf3, "Todos se leen con st_read() del paquete sf", font_size=14, color=GRAY_TEXT)

slide_geom_sf()


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 5: mutate()
# ═══════════════════════════════════════════════════════════════════════════
add_function_slide(
    func_name="mutate()",
    package_name="dplyr",
    description="Crea nuevas columnas o modifica columnas existentes en un dataframe. "
                "Es uno de los verbos más usados del tidyverse. Las columnas nuevas "
                "se agregan al final del dataframe.",
    syntax="datos %>%\n  mutate(nueva_col = expresion)",
    example_code=(
        "library(tidyverse)\n\n"
        "# Crear nuevas columnas\n"
        "mtcars %>%\n"
        "  mutate(\n"
        "    kpl = mpg * 0.425,          # millas a km/litro\n"
        "    peso_ton = wt / 2.205,      # libras a toneladas\n"
        "    potencia_cat = ifelse(\n"
        '      hp > 150, "Alta", "Baja"\n'
        "    )\n"
        "  )"
    ),
    notes=[
        "Siempre devuelve un dataframe del mismo número de filas",
        "Puedes crear múltiples columnas en un solo mutate()",
        "Las columnas nuevas pueden usar columnas recién creadas",
        "Si el nombre ya existe, sobreescribe la columna",
        "Dentro de mutate puedes usar: ifelse(), case_when(),\n  str_c(), round(), y muchas más funciones",
        "Para modificar solo un subconjunto, mira mutate() + if_else()",
    ],
    accent_color=ACCENT_PEACH,
    tip="Usa .before o .after para controlar dónde se inserta la nueva columna:\n"
        "mutate(nueva = x + y, .after = columna_ref)"
)


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 6: case_when() e ifelse()
# ═══════════════════════════════════════════════════════════════════════════
def slide_case_when_ifelse():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 2.0, 0.45, "dplyr", ACCENT_MAUVE)
    add_textbox(slide, 0.5, 1.0, 14, 0.8, "case_when() e ifelse()",
                font_size=40, color=ACCENT_MAUVE, bold=True, font_name="Ubuntu Mono")

    tf = add_rich_textbox(slide, 0.5, 1.9, 7.0, 0.8)
    add_paragraph(tf, "Funciones para crear valores condicionales dentro de mutate().",
                  font_size=18, color=GRAY_TEXT)

    # ifelse
    add_textbox(slide, 0.5, 2.8, 5, 0.4, "ifelse() — Dos opciones",
                font_size=20, color=ACCENT_PEACH, bold=True)
    add_code_block(slide, 0.5, 3.2, 7.0, 1.8,
        "# ifelse: condicion -> valor_si_TRUE, valor_si_FALSE\n"
        "datos %>%\n"
        "  mutate(\n"
        "    categoria = ifelse(\n"
        "      precio > 100,\n"
        '      "Caro",\n'
        '      "Barato"\n'
        "    )\n"
        "  )", font_size=13)

    # case_when
    add_textbox(slide, 0.5, 5.2, 5, 0.4, "case_when() — Múltiples opciones",
                font_size=20, color=ACCENT_GREEN, bold=True)
    add_code_block(slide, 0.5, 5.6, 7.0, 2.8,
        "# case_when: multiples condiciones con ~\n"
        "datos %>%\n"
        "  mutate(\n"
        "    nivel = case_when(\n"
        '      calificacion >= 90 ~ "Excelente",\n'
        '      calificacion >= 80 ~ "Bueno",\n'
        '      calificacion >= 70 ~ "Regular",\n'
        '      TRUE               ~ "Insuficiente"\n'
        '      # TRUE es el "else" (todo lo demas)\n'
        "    )\n"
        "  )", font_size=13)

    # Panel derecho
    add_textbox(slide, 8.5, 1.0, 7, 0.5, "Comparación",
                font_size=22, color=DARK_TEXT, bold=True)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.03))
    sep.fill.solid(); sep.fill.fore_color.rgb = ACCENT_MAUVE; sep.line.fill.background()

    tf2 = add_rich_textbox(slide, 8.5, 1.7, 6.8, 3.5)
    add_paragraph(tf2, "ifelse()", font_size=18, color=ACCENT_PEACH, bold=True, font_name="Ubuntu Mono")
    add_paragraph(tf2, "  + Simple: solo TRUE / FALSE", font_size=15, color=GRAY_TEXT)
    add_paragraph(tf2, "  + Rápido para decisiones binarias", font_size=15, color=GRAY_TEXT)
    add_paragraph(tf2, "  - Se vuelve feo si anidas muchos", font_size=15, color=ACCENT_RED)
    add_paragraph(tf2, "", font_size=10, color=GRAY_TEXT)
    add_paragraph(tf2, "case_when()", font_size=18, color=ACCENT_GREEN, bold=True, font_name="Ubuntu Mono")
    add_paragraph(tf2, "  + Múltiples condiciones, limpio", font_size=15, color=GRAY_TEXT)
    add_paragraph(tf2, "  + Más legible que ifelse anidados", font_size=15, color=GRAY_TEXT)
    add_paragraph(tf2, '  + TRUE ~ captura "todo lo demás"', font_size=15, color=GRAY_TEXT)
    add_paragraph(tf2, "  - Sintaxis con ~ puede confundir al inicio", font_size=15, color=ACCENT_RED)

    # Cuándo usar cuál
    add_textbox(slide, 8.5, 5.5, 7, 0.4, "¿Cuándo usar cuál?",
                font_size=18, color=DARK_TEXT, bold=True)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(5.9), Inches(6.5), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = ACCENT_MAUVE; sep2.line.fill.background()

    tf3 = add_rich_textbox(slide, 8.5, 6.1, 6.8, 2.0)
    add_paragraph(tf3, "2 opciones  -->  ifelse()", font_size=16, color=ACCENT_PEACH, bold=True)
    add_paragraph(tf3, "3+ opciones -->  case_when()", font_size=16, color=ACCENT_GREEN, bold=True)
    add_paragraph(tf3, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf3, "Tip: También existe if_else() de dplyr, que es "
                  "más estricto con los tipos de datos (recomendado).",
                  font_size=14, color=GRAY_TEXT)

slide_case_when_ifelse()


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 7: janitor::clean_names()
# ═══════════════════════════════════════════════════════════════════════════
add_function_slide(
    func_name="clean_names()",
    package_name="janitor",
    description="Limpia los nombres de las columnas de un dataframe: "
                "los convierte a snake_case, elimina caracteres especiales, "
                "acentos, espacios y los hace consistentes.",
    syntax="library(janitor)\n\ndatos %>%\n  clean_names()",
    example_code=(
        "library(tidyverse)\n"
        "library(janitor)\n\n"
        "# Datos con nombres feos\n"
        "datos <- tibble(\n"
        '  `Nombre Completo` = c("Ana", "Luis"),\n'
        "  `Año de Nacimiento` = c(1995, 1998),\n"
        "  `% Asistencia` = c(95.5, 88.2)\n"
        ")\n\n"
        "datos %>% clean_names()\n"
        "# nombre_completo, ano_de_nacimiento,\n"
        "# percent_asistencia"
    ),
    notes=[
        "Convierte TODO a snake_case (minúsculas con _)",
        "Elimina acentos: Año -> ano",
        "Reemplaza espacios por guiones bajos",
        "Convierte % a 'percent', # a 'number'",
        "Elimina caracteres especiales",
        "Si hay nombres duplicados, agrega sufijo numérico",
        "Funciona con dataframes, tibbles y más",
        "Úsalo siempre al importar datos externos",
    ],
    accent_color=ACCENT_YELLOW,
    tip="Ponlo justo después de leer tus datos:\n"
        "read_csv('archivo.csv') %>% clean_names()"
)


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 8: str_c()
# ═══════════════════════════════════════════════════════════════════════════
add_function_slide(
    func_name="str_c()",
    package_name="stringr",
    description="Concatena (pega) cadenas de texto. Es la versión tidyverse de paste0(). "
                "Más predecible con valores NA que paste().",
    syntax='str_c("texto1", "texto2", sep = " ")',
    example_code=(
        "library(tidyverse)\n\n"
        "# Concatenar columnas\n"
        "datos <- tibble(\n"
        '  nombre = c("Ana", "Luis"),\n'
        '  apellido = c("García", "López")\n'
        ")\n\n"
        "datos %>%\n"
        "  mutate(\n"
        "    nombre_completo = str_c(nombre, apellido,\n"
        '                            sep = " "),\n'
        '    saludo = str_c("Hola, ", nombre, "!")\n'
        "  )\n"
        '# "Ana García", "Hola, Ana!"'
    ),
    notes=[
        "sep = define el separador entre elementos",
        "collapse = une un vector en un solo string",
        "Si algún valor es NA, el resultado es NA\n  (a diferencia de paste que pone \"NA\")",
        "str_c('a', 'b')        -> 'ab'  (sin separador)",
        "str_c('a', 'b', sep='-') -> 'a-b'",
        "Equivalente a paste0() pero más consistente",
    ],
    accent_color=ACCENT_BLUE,
    tip='Usa collapse para colapsar un vector:\nstr_c(c("a","b","c"), collapse=", ") -> "a, b, c"'
)


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 9: round()
# ═══════════════════════════════════════════════════════════════════════════
add_function_slide(
    func_name="round()",
    package_name="base R",
    description="Redondea números al número de decimales especificado. "
                "Función base de R, no requiere paquetes adicionales.",
    syntax="round(x, digits = 0)",
    example_code=(
        "# Redondeo básico\n"
        "round(3.14159, 2)     # 3.14\n"
        "round(3.14159, 0)     # 3\n"
        "round(1234.5, -2)     # 1200\n\n"
        "# Dentro de mutate\n"
        "datos %>%\n"
        "  mutate(\n"
        "    promedio = round(mean(valor), 2),\n"
        "    porcentaje = round(n / total * 100, 1)\n"
        "  )"
    ),
    notes=[
        "digits = 0 -> sin decimales (entero)",
        "digits = 2 -> dos decimales",
        "digits negativo -> redondea a decenas, centenas, etc.",
        "round(1250, -2) -> 1200",
        "Funciones relacionadas:\n"
        "  ceiling() -> redondea hacia arriba\n"
        "  floor()   -> redondea hacia abajo\n"
        "  trunc()   -> elimina decimales",
        "R usa 'banker's rounding': round(2.5) = 2\n"
        "  (redondea al par más cercano en empates)",
    ],
    accent_color=ACCENT_GREEN,
    tip="Si necesitas siempre redondear .5 hacia arriba:\n"
        "round2 <- function(x, n=0) floor(x*10^n + 0.5)/10^n"
)


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 10: ungroup()
# ═══════════════════════════════════════════════════════════════════════════
add_function_slide(
    func_name="ungroup()",
    package_name="dplyr",
    description="Elimina la agrupación de un dataframe creada con group_by(). "
                "Es fundamental para evitar cálculos inesperados en operaciones "
                "posteriores.",
    syntax="datos %>%\n  group_by(variable) %>%\n  summarise(...) %>%\n  ungroup()",
    example_code=(
        "library(tidyverse)\n\n"
        "# Sin ungroup: el mutate opera por grupo\n"
        "mtcars %>%\n"
        "  group_by(cyl) %>%\n"
        "  summarise(promedio = mean(mpg)) %>%\n"
        "  ungroup() %>%    # <- Importante!\n"
        "  mutate(\n"
        "    porcentaje = promedio / sum(promedio) * 100\n"
        "  )\n\n"
        "# Sin ungroup, sum(promedio) seria por grupo\n"
        "# Con ungroup, sum(promedio) es el total global"
    ),
    notes=[
        "group_by() se 'hereda' en todas las operaciones\n  posteriores hasta que hagas ungroup()",
        "summarise() quita UN nivel de agrupación,\n  pero si tenías group_by(a, b), queda agrupado por 'a'",
        "Buena práctica: siempre usar ungroup() después\n  de terminar operaciones agrupadas",
        "Si no desagrupas, mutate(), filter() y otras\n  funciones operarán dentro de cada grupo",
        "Puedes verificar con: groups(datos)",
    ],
    accent_color=ACCENT_RED,
    tip="Regla de oro: si usas group_by(), siempre\ntermina con ungroup() cuando ya no lo necesites."
)


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 11: aes() — dentro vs fuera
# ═══════════════════════════════════════════════════════════════════════════
def slide_aes_dentro_fuera():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 2.0, 0.45, "ggplot2", BLUE_TEC)
    add_textbox(slide, 0.5, 1.0, 14, 0.8, "aes() — Dentro vs Fuera del geom",
                font_size=36, color=DARK_TEXT, bold=True)

    tf_intro = add_rich_textbox(slide, 0.5, 1.9, 14.5, 0.6)
    add_paragraph(tf_intro,
        "La ubicación de los parámetros estéticos cambia completamente el comportamiento de la gráfica.",
        font_size=18, color=GRAY_TEXT)

    # --- Columna izquierda: DENTRO de aes() ---
    add_textbox(slide, 0.5, 2.8, 7, 0.5, "Dentro de aes() — Mapeo a datos",
                font_size=22, color=ACCENT_GREEN, bold=True)
    sep1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(3.3), Inches(7.0), Inches(0.03))
    sep1.fill.solid(); sep1.fill.fore_color.rgb = ACCENT_GREEN; sep1.line.fill.background()

    add_code_block(slide, 0.5, 3.5, 7.0, 1.8,
        "# El COLOR depende de una variable\n"
        "ggplot(datos, aes(x = mes, y = ventas)) +\n"
        "  geom_col(aes(fill = region))\n\n"
        "# Cada región tendrá un color diferente\n"
        "# ggplot asigna colores automáticamente",
        font_size=13)

    tf_left = add_rich_textbox(slide, 0.5, 5.5, 7.0, 2.8)
    add_paragraph(tf_left, "  El valor cambia según los datos", font_size=16, color=ACCENT_GREEN)
    add_paragraph(tf_left, "  Se genera una leyenda automáticamente", font_size=16, color=ACCENT_GREEN)
    add_paragraph(tf_left, "  Usa nombres de columnas (sin comillas)", font_size=16, color=ACCENT_GREEN)
    add_paragraph(tf_left, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf_left, "  Ejemplo: aes(fill = region)", font_size=16, color=GRAY_TEXT, font_name="Ubuntu Mono")
    add_paragraph(tf_left, "  -> Norte=azul, Sur=rojo, Centro=verde...", font_size=15, color=GRAY_TEXT)

    # --- Columna derecha: FUERA de aes() ---
    add_textbox(slide, 8.5, 2.8, 7, 0.5, "Fuera de aes() — Valor fijo",
                font_size=22, color=ACCENT_PEACH, bold=True)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(3.3), Inches(6.5), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = ACCENT_PEACH; sep2.line.fill.background()

    add_code_block(slide, 8.5, 3.5, 6.8, 1.8,
        "# TODAS las barras del mismo color\n"
        "ggplot(datos, aes(x = mes, y = ventas)) +\n"
        '  geom_col(fill = "#1A6BD4")\n\n'
        "# Un solo color fijo para TODO el geom\n"
        "# No se genera leyenda",
        font_size=13)

    tf_right = add_rich_textbox(slide, 8.5, 5.5, 6.8, 2.8)
    add_paragraph(tf_right, "  Aplica el MISMO valor a todos los elementos", font_size=16, color=ACCENT_PEACH)
    add_paragraph(tf_right, "  No genera leyenda", font_size=16, color=ACCENT_PEACH)
    add_paragraph(tf_right, "  Usa valores literales (con comillas para colores)", font_size=16, color=ACCENT_PEACH)
    add_paragraph(tf_right, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf_right, '  Ejemplo: fill = "#1A6BD4"', font_size=16, color=GRAY_TEXT, font_name="Ubuntu Mono")
    add_paragraph(tf_right, "  -> Todas las barras: azul", font_size=15, color=GRAY_TEXT)

    # Línea central divisoria
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(7.85), Inches(2.8), Inches(0.03), Inches(5.3))
    div.fill.solid(); div.fill.fore_color.rgb = LIGHT_GRAY; div.line.fill.background()

slide_aes_dentro_fuera()


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 12: prettyNum()
# ═══════════════════════════════════════════════════════════════════════════
add_function_slide(
    func_name="prettyNum()",
    package_name="base R",
    description="Formatea números para presentación: agrega separadores de miles, "
                "controla decimales y devuelve texto legible. "
                "Ideal para etiquetas en gráficas y tablas.",
    syntax='prettyNum(x, big.mark = ",", scientific = FALSE)',
    example_code=(
        "# Formato con comas\n"
        'prettyNum(1234567, big.mark = ",")\n'
        '# "1,234,567"\n\n'
        "# Dentro de mutate para etiquetas\n"
        "datos %>%\n"
        "  mutate(\n"
        "    etiqueta = prettyNum(\n"
        "      poblacion,\n"
        '      big.mark = ",",\n'
        "      scientific = FALSE\n"
        "    )\n"
        "  )\n\n"
        "# En gráficas\n"
        "geom_text(aes(label = prettyNum(valor,\n"
        '               big.mark = \",\")))'
    ),
    notes=[
        'big.mark = "," -> separador de miles con coma',
        'big.mark = "." -> separador con punto (estilo MX)',
        "scientific = FALSE -> evita notación científica",
        "Devuelve un CHARACTER, no un número",
        "Muy útil para geom_text() y geom_label()",
        "Para más control, revisa scales::comma()\n  y format()",
    ],
    accent_color=ACCENT_MAUVE,
    tip='Para formato en español/mexicano:\nprettyNum(1234567, big.mark = ".") -> "1.234.567"'
)


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 13: Medidor de color digital
# ═══════════════════════════════════════════════════════════════════════════
def slide_color_picker():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 3.0, 0.45, "Herramienta", ACCENT_PINK)
    add_textbox(slide, 0.5, 1.0, 14, 0.8, "Medidor de color digital (Color Picker)",
                font_size=36, color=DARK_TEXT, bold=True)

    tf = add_rich_textbox(slide, 0.5, 1.9, 14.5, 0.8)
    add_paragraph(tf,
        "Herramienta para obtener el código hexadecimal exacto de cualquier color en pantalla.",
        font_size=18, color=GRAY_TEXT)

    # ¿Qué es un código hex?
    add_textbox(slide, 0.5, 2.9, 7, 0.4, "¿Qué es un código hexadecimal?",
                font_size=22, color=DARK_TEXT, bold=True)
    sep1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(3.3), Inches(7.0), Inches(0.03))
    sep1.fill.solid(); sep1.fill.fore_color.rgb = ACCENT_PINK; sep1.line.fill.background()

    tf2 = add_rich_textbox(slide, 0.5, 3.5, 7.0, 1.5)
    add_paragraph(tf2, "Formato: #RRGGBB", font_size=20, color=ACCENT_PINK, bold=True, font_name="Ubuntu Mono")
    add_paragraph(tf2, "Cada par de letras = intensidad de Rojo, Verde, Azul",
                  font_size=16, color=GRAY_TEXT)
    add_paragraph(tf2, "Rango: 00 (nada) a FF (máximo)", font_size=16, color=GRAY_TEXT)

    # Ejemplos de colores
    colors_examples = [
        ("#FF0000", "Rojo puro", RGBColor(0xCC, 0x00, 0x00)),
        ("#00AA00", "Verde", RGBColor(0x00, 0xAA, 0x00)),
        ("#0000FF", "Azul puro", RGBColor(0x00, 0x00, 0xCC)),
        ("#003DA5", "Azul Tec", BLUE_TEC),
        ("#C0397A", "Rosa", ACCENT_PINK),
    ]
    y_pos = 5.2
    for hex_code, name, color in colors_examples:
        # Cuadro de color
        swatch = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos), Inches(0.5), Inches(0.4)
        )
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = RGBColor(
            int(hex_code[1:3], 16), int(hex_code[3:5], 16), int(hex_code[5:7], 16)
        )
        swatch.line.color.rgb = LIGHT_GRAY
        swatch.line.width = Pt(1)
        swatch.adjustments[0] = 0.1

        add_textbox(slide, 1.2, y_pos, 2.5, 0.4, hex_code,
                    font_size=16, color=color, bold=True, font_name="Ubuntu Mono")
        add_textbox(slide, 3.5, y_pos, 3, 0.4, name,
                    font_size=14, color=GRAY_TEXT)
        y_pos += 0.5

    # Panel derecho: Herramientas
    add_textbox(slide, 8.5, 1.0, 7, 0.5, "¿Cómo obtener colores?",
                font_size=22, color=DARK_TEXT, bold=True)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = ACCENT_PINK; sep2.line.fill.background()

    tf3 = add_rich_textbox(slide, 8.5, 1.7, 6.8, 2.5)
    add_paragraph(tf3, "En macOS:", font_size=18, color=ACCENT_TEAL, bold=True)
    add_paragraph(tf3, "  Utilidades -> Medidor de color digital", font_size=15, color=GRAY_TEXT)
    add_paragraph(tf3, '  O buscar "Digital Color Meter" en Spotlight', font_size=15, color=GRAY_TEXT)
    add_paragraph(tf3, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf3, "En Windows:", font_size=18, color=ACCENT_TEAL, bold=True)
    add_paragraph(tf3, "  PowerToys -> Color Picker (Win+Shift+C)", font_size=15, color=GRAY_TEXT)
    add_paragraph(tf3, "", font_size=8, color=GRAY_TEXT)
    add_paragraph(tf3, "En la web:", font_size=18, color=ACCENT_TEAL, bold=True)
    add_paragraph(tf3, '  Google: buscar "color picker"', font_size=15, color=GRAY_TEXT)
    add_paragraph(tf3, "  coolors.co — paletas de colores", font_size=15, color=GRAY_TEXT)

    # Uso en R
    add_textbox(slide, 8.5, 4.8, 7, 0.4, "Uso en ggplot2",
                font_size=22, color=DARK_TEXT, bold=True)
    sep3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(5.2), Inches(6.5), Inches(0.03))
    sep3.fill.solid(); sep3.fill.fore_color.rgb = ACCENT_PINK; sep3.line.fill.background()

    add_code_block(slide, 8.5, 5.4, 6.8, 2.8,
        "# Usar el código hex directamente\n"
        'geom_col(fill = "#003DA5")\n'
        'geom_point(color = "#C0397A")\n\n'
        "# Para múltiples colores manuales\n"
        "scale_fill_manual(\n"
        '  values = c("#003DA5", "#1B8A4A",\n'
        '             "#C0397A", "#D06B21")\n'
        ")",
        font_size=13)

slide_color_picker()


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 14: ggsave()
# ═══════════════════════════════════════════════════════════════════════════
add_function_slide(
    func_name="ggsave()",
    package_name="ggplot2",
    description="Guarda la última gráfica creada con ggplot2 en un archivo. "
                "Detecta el formato automáticamente por la extensión del archivo. "
                "Permite controlar tamaño, resolución y más.",
    syntax='ggsave("nombre.png", width = 10, height = 6, dpi = 300)',
    example_code=(
        "library(tidyverse)\n\n"
        "# Crear gráfica\n"
        "mi_grafica <- ggplot(mpg, aes(displ, hwy)) +\n"
        "  geom_point()\n\n"
        "# Guardar como PNG (alta resolución)\n"
        'ggsave("grafica.png",\n'
        "       plot = mi_grafica,\n"
        "       width = 10, height = 6,\n"
        "       dpi = 300)\n\n"
        "# Guardar como SVG (vectorial)\n"
        'ggsave("grafica.svg",\n'
        "       width = 10, height = 6)"
    ),
    notes=[
        "Si no especificas plot =, guarda la ÚLTIMA gráfica",
        "width y height están en pulgadas por defecto",
        "units = 'cm' o 'mm' para cambiar unidades",
        "dpi = 300 es estándar para impresión",
        "dpi = 72 es suficiente para pantalla/web",
        "Formatos soportados: png, jpeg, svg, pdf, tiff, bmp",
        "Para guardar en carpeta:\n"
        '  ggsave("carpeta/grafica.png")',
    ],
    accent_color=ACCENT_PEACH,
    tip="Para presentaciones usa width=16, height=9 (16:9).\n"
        "Para redes sociales: width=10, height=10 (cuadrado)."
)


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 15: PNG vs JPEG vs SVG
# ═══════════════════════════════════════════════════════════════════════════
def slide_formats():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 3.5, 0.45, "Formatos de imagen", ACCENT_TEAL)
    add_textbox(slide, 0.5, 1.0, 14, 0.8, "PNG vs JPEG vs SVG",
                font_size=40, color=DARK_TEXT, bold=True)

    tf_intro = add_rich_textbox(slide, 0.5, 1.9, 14.5, 0.6)
    add_paragraph(tf_intro,
        "Cada formato tiene ventajas y desventajas. Elegir el correcto depende del uso.",
        font_size=18, color=GRAY_TEXT)

    # ── PNG ──
    col_x = 0.5
    add_textbox(slide, col_x, 2.8, 4.5, 0.5, "PNG",
                font_size=28, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, col_x, 3.3, 4.5, 0.35, "Portable Network Graphics",
                font_size=13, color=GRAY_TEXT)
    sep1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(col_x), Inches(3.65), Inches(4.3), Inches(0.03))
    sep1.fill.solid(); sep1.fill.fore_color.rgb = ACCENT_BLUE; sep1.line.fill.background()

    tf_png = add_rich_textbox(slide, col_x, 3.8, 4.3, 4.5)
    add_paragraph(tf_png, "Tipo: Raster (píxeles)", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_png, "Compresión: Sin pérdida", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_png, "Transparencia: Sí (canal alpha)", font_size=14, color=ACCENT_GREEN)
    add_paragraph(tf_png, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_png, "Ventajas:", font_size=15, color=ACCENT_BLUE, bold=True)
    add_paragraph(tf_png, "  Calidad perfecta, sin artefactos", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_png, "  Soporta transparencia", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_png, "  Ideal para gráficas y diagramas", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_png, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_png, "Desventajas:", font_size=15, color=ACCENT_RED, bold=True)
    add_paragraph(tf_png, "  Archivos más pesados que JPEG", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_png, "  Pierde calidad al escalar", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_png, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_png, "Usar para:", font_size=15, color=DARK_TEXT, bold=True)
    add_paragraph(tf_png, "  Gráficas, capturas, logos, web", font_size=14, color=ACCENT_BLUE, bold=True)

    # ── JPEG ──
    col_x = 5.5
    add_textbox(slide, col_x, 2.8, 4.5, 0.5, "JPEG",
                font_size=28, color=ACCENT_PEACH, bold=True)
    add_textbox(slide, col_x, 3.3, 4.5, 0.35, "Joint Photographic Experts Group",
                font_size=13, color=GRAY_TEXT)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(col_x), Inches(3.65), Inches(4.3), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = ACCENT_PEACH; sep2.line.fill.background()

    tf_jpg = add_rich_textbox(slide, col_x, 3.8, 4.3, 4.5)
    add_paragraph(tf_jpg, "Tipo: Raster (píxeles)", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "Compresión: Con pérdida", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "Transparencia: No soporta", font_size=14, color=ACCENT_RED)
    add_paragraph(tf_jpg, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "Ventajas:", font_size=15, color=ACCENT_PEACH, bold=True)
    add_paragraph(tf_jpg, "  Archivos muy ligeros", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "  Excelente para fotografías", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "  Compatible con todo", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "Desventajas:", font_size=15, color=ACCENT_RED, bold=True)
    add_paragraph(tf_jpg, "  Pierde calidad (artefactos)", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "  Sin transparencia", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_jpg, "Usar para:", font_size=15, color=DARK_TEXT, bold=True)
    add_paragraph(tf_jpg, "  Fotos, imágenes con muchos\n  colores, web si el peso importa", font_size=14, color=ACCENT_PEACH, bold=True)

    # ── SVG ──
    col_x = 10.5
    add_textbox(slide, col_x, 2.8, 4.5, 0.5, "SVG",
                font_size=28, color=ACCENT_GREEN, bold=True)
    add_textbox(slide, col_x, 3.3, 4.5, 0.35, "Scalable Vector Graphics",
                font_size=13, color=GRAY_TEXT)
    sep3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(col_x), Inches(3.65), Inches(4.3), Inches(0.03))
    sep3.fill.solid(); sep3.fill.fore_color.rgb = ACCENT_GREEN; sep3.line.fill.background()

    tf_svg = add_rich_textbox(slide, col_x, 3.8, 4.3, 4.5)
    add_paragraph(tf_svg, "Tipo: Vectorial (matemáticas)", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_svg, "Compresión: No aplica", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_svg, "Transparencia: Sí", font_size=14, color=ACCENT_GREEN)
    add_paragraph(tf_svg, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_svg, "Ventajas:", font_size=15, color=ACCENT_GREEN, bold=True)
    add_paragraph(tf_svg, "  Escala infinitamente sin perder\n    calidad", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_svg, "  Archivos pequeños para gráficas", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_svg, "  Editable (es código XML)", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_svg, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_svg, "Desventajas:", font_size=15, color=ACCENT_RED, bold=True)
    add_paragraph(tf_svg, "  No ideal para fotos", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_svg, "  Puede ser pesado si es complejo", font_size=14, color=GRAY_TEXT)
    add_paragraph(tf_svg, "", font_size=6, color=GRAY_TEXT)
    add_paragraph(tf_svg, "Usar para:", font_size=15, color=DARK_TEXT, bold=True)
    add_paragraph(tf_svg, "  Gráficas para publicación,\n  logos, íconos, presentaciones", font_size=14, color=ACCENT_GREEN, bold=True)

    # Líneas divisorias verticales
    for x in [5.15, 10.15]:
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2.8), Inches(0.02), Inches(5.5))
        div.fill.solid(); div.fill.fore_color.rgb = LIGHT_GRAY; div.line.fill.background()

    # Recomendación al pie
    rec_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(8.0), Inches(15.0), Inches(0.5)
    )
    rec_shape.fill.solid()
    rec_shape.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)  # azul muy pastel
    rec_shape.line.color.rgb = ACCENT_BLUE
    rec_shape.line.width = Pt(1)
    rec_shape.adjustments[0] = 0.15
    tf_rec = rec_shape.text_frame
    tf_rec.word_wrap = True
    tf_rec.margin_left = Inches(0.3)
    p = tf_rec.paragraphs[0]
    p.text = "Recomendación para gráficas en R:  PNG (uso general, dpi=300)  |  SVG (publicaciones, máxima calidad)  |  JPEG (solo si necesitas archivos muy ligeros)"
    p.font.size = Pt(15)
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Ubuntu"
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

slide_formats()


# ═══════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 16: Vectores con nombre (Named Vectors)
# ═══════════════════════════════════════════════════════════════════════════
def slide_named_vectors():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_WHITE)
    add_footer_bar(slide)

    add_badge(slide, 0.5, 0.4, 2.0, 0.45, "base R", BLUE_TEC)
    add_textbox(slide, 0.5, 1.0, 10, 0.8, "Named Vectors (Vectores con nombre)",
                font_size=36, color=DARK_TEXT, bold=True)

    tf = add_rich_textbox(slide, 0.5, 1.9, 7.0, 1.4)
    add_paragraph(tf, "Un named vector es un vector donde cada elemento tiene "
                  "un nombre asociado. Esto permite acceder a valores por nombre "
                  "en lugar de por posición, y es clave para mapear valores en ggplot2.",
                  font_size=18, color=GRAY_TEXT)

    add_textbox(slide, 0.5, 3.3, 3, 0.4, "Sintaxis:", font_size=14,
                color=BLUE_TEC, bold=True)
    add_code_block(slide, 0.5, 3.7, 7.0, 1.0,
        'mi_vector <- c("nombre_1" = valor1,\n'
        '               "nombre_2" = valor2,\n'
        '               "nombre_3" = valor3)', font_size=14)

    add_textbox(slide, 0.5, 4.9, 3, 0.4, "Ejemplos:", font_size=14,
                color=BLUE_TEC, bold=True)
    add_code_block(slide, 0.5, 5.3, 7.0, 3.0,
        "# Definir colores para categorías\n"
        'colores <- c("Norte"  = "#003DA5",\n'
        '             "Centro" = "#1B8A4A",\n'
        '             "Sur"    = "#C0397A")\n\n'
        "# Usar en ggplot con scale_*_manual()\n"
        "ggplot(datos, aes(x = region, y = valor,\n"
        "                  fill = region)) +\n"
        "  geom_col() +\n"
        "  scale_fill_manual(values = colores)\n\n"
        "# Acceder por nombre\n"
        'colores["Norte"]   # "#003DA5"',
        font_size=12)

    # Panel derecho
    add_textbox(slide, 8.5, 1.0, 7, 0.5, "¿Para qué sirven?",
                font_size=22, color=DARK_TEXT, bold=True)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.03))
    sep.fill.solid(); sep.fill.fore_color.rgb = BLUE_TEC; sep.line.fill.background()

    tf2 = add_rich_textbox(slide, 8.5, 1.7, 6.8, 2.8)
    uses = [
        "Asignar colores fijos a categorías en\n  scale_fill_manual() / scale_color_manual()",
        "Renombrar o recodificar valores con\n  recode() o deframe()",
        "Acceder a valores por nombre en vez de\n  por posición numérica",
        "Crear etiquetas personalizadas para facetas\n  o ejes en gráficas",
        "Definir parámetros de configuración\n  reutilizables en tu script",
    ]
    for use in uses:
        add_paragraph(tf2, f"  {use}", font_size=15, color=GRAY_TEXT, space_after=Pt(10))

    # Ejemplo scale_fill_manual
    add_textbox(slide, 8.5, 5.0, 7, 0.4, "Uso típico en ggplot2",
                font_size=18, color=DARK_TEXT, bold=True)
    sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(5.4), Inches(6.5), Inches(0.03))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = BLUE_TEC; sep2.line.fill.background()

    tf3 = add_rich_textbox(slide, 8.5, 5.6, 6.8, 1.2)
    add_paragraph(tf3, "El nombre del vector debe coincidir exactamente",
                  font_size=15, color=GRAY_TEXT)
    add_paragraph(tf3, "con los valores de la variable en los datos:",
                  font_size=15, color=GRAY_TEXT)

    add_code_block(slide, 8.5, 6.9, 6.8, 1.3,
        "# Si la columna 'region' tiene: Norte, Centro, Sur\n"
        "# Los nombres del vector deben ser iguales:\n"
        'c("Norte" = "#003DA5", "Centro" = "#1B8A4A",\n'
        '  "Sur" = "#C0397A")',
        font_size=12)

slide_named_vectors()


# ═══════════════════════════════════════════════════════════════════════════
# Guardar presentación
# ═══════════════════════════════════════════════════════════════════════════
output_path = "/Users/jorgejuvenalcamposferreira/Documents/GitHub/TC2001B.601-Ciencia-de-datos-ene-jun-2026/01_PRESENTACIONES/Sesión_07/Fichas_Sesion_07.pptx"
prs.save(output_path)
print(f"Presentación guardada en: {output_path}")
print(f"Total de diapositivas: {len(prs.slides)}")
