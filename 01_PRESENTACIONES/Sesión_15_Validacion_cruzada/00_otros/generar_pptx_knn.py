"""Construye la presentacion PPTX de K-Vecinos Cercanos estilo Tec."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === Paleta Tec ===
NAVY = RGBColor(0x1E, 0x4C, 0x7D)
LIGHT_BLUE = RGBColor(0xEA, 0xF2, 0xFA)
SOFT_BLUE = RGBColor(0xDD, 0xE5, 0xED)
SOFT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
DARK_GREEN = RGBColor(0x1A, 0x7A, 0x3A)
SOFT_RED = RGBColor(0xFA, 0xDB, 0xD8)
DARK_RED = RGBColor(0xC0, 0x39, 0x2B)
HIGHLIGHT_GREEN = RGBColor(0xF0, 0xFA, 0xF0)
GRAY_TOP = RGBColor(0xF0, 0xF0, 0xF0)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
UBUNTU = 'Ubuntu'

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

TITULO_CORTO = "K-Vecinos Cercanos"


def add_text(slide, text, left, top, width, height, *, size=14, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, italic=False, line_spacing=None,
             font=UBUNTU, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_rect(slide, left, top, width, height, fill_color, *, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_bg(slide, color=WHITE):
    bg = add_rect(slide, 0, 0, SW, SH, color)
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_header(slide, title, *, width_in=6.5):
    add_rect(slide, 0, Inches(0.3), Inches(width_in), Inches(0.7), NAVY)
    add_text(slide, title, Inches(0.0), Inches(0.3), Inches(width_in),
             Inches(0.7), size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_card(slide, left, top, width, height, *,
             fill=LIGHT_BLUE, title=None, body=None,
             title_color=NAVY, body_color=BLACK,
             title_size=14, body_size=12, line_spacing=1.3):
    add_rect(slide, left, top, width, height, fill)
    if title and body:
        add_text(slide, title, left + Inches(0.15), top + Inches(0.1),
                 width - Inches(0.3), Inches(0.4), size=title_size,
                 bold=True, color=title_color)
        add_text(slide, body, left + Inches(0.15), top + Inches(0.55),
                 width - Inches(0.3), height - Inches(0.65),
                 size=body_size, color=body_color, line_spacing=line_spacing)
    elif title:
        add_text(slide, title, left + Inches(0.15), top + Inches(0.1),
                 width - Inches(0.3), height - Inches(0.2),
                 size=title_size, bold=True, color=title_color,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=line_spacing)
    elif body:
        add_text(slide, body, left + Inches(0.15), top + Inches(0.12),
                 width - Inches(0.3), height - Inches(0.24),
                 size=body_size, color=body_color, line_spacing=line_spacing)


def add_formula_box(slide, formula, sub, left, top, width, height):
    add_rect(slide, left, top, width, height, NAVY)
    add_text(slide, formula, left, top + Inches(0.18),
             width, Inches(0.6), size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, left, top + height - Inches(0.4),
                 width, Inches(0.35), size=12, color=WHITE,
                 align=PP_ALIGN.CENTER, italic=True)


def add_footer(slide, num):
    add_text(slide, f"{TITULO_CORTO} - {num}",
             Inches(7.0), Inches(7.2), Inches(2.8), Inches(0.2),
             size=9, color=GRAY, align=PP_ALIGN.RIGHT, italic=True)


# ============ SLIDES ============

# --- Slide 1: Portada ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, SW, Inches(1.1), GRAY_TOP)
add_text(s, "Tecnologico de Monterrey", Inches(0.4), Inches(0.25),
         Inches(6), Inches(0.6), size=20, color=NAVY, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "TC2001B  Ciencia de datos", Inches(5.5), Inches(0.25),
         Inches(4.1), Inches(0.6), size=12, color=GRAY, italic=True,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, 0, Inches(1.1), SW, Inches(3.7), SOFT_BLUE)
add_text(s, "K-Vecinos Cercanos", Inches(0.6), Inches(1.7),
         Inches(8.8), Inches(1.3), size=54, color=NAVY, bold=True)
add_text(s, "(K-NN)", Inches(0.6), Inches(2.95),
         Inches(8.8), Inches(0.8), size=40, color=DARK_RED, bold=True)
add_text(s, "Un algoritmo intuitivo para clasificacion y regresion",
         Inches(0.6), Inches(3.95), Inches(8.8), Inches(0.6),
         size=18, color=NAVY, italic=True)

add_rect(s, Inches(0.6), Inches(6.0), Inches(0.05), Inches(0.6), NAVY)
add_text(s, "Jorge Juvenal Campos Ferreira", Inches(0.8), Inches(6.0),
         Inches(7), Inches(0.35), size=13, color=BLACK, bold=True)
add_rect(s, Inches(0.8), Inches(6.45), Inches(0.15), Inches(0.15), DARK_RED)
add_text(s, "juvenal.campos@tec.mx", Inches(1.05), Inches(6.4),
         Inches(7), Inches(0.3), size=11, color=DARK_RED)


# --- Slide 2: Que es K-NN ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Que es K-NN?")
add_text(s,
         "Algoritmo de aprendizaje supervisado propuesto formalmente por "
         "Fix y Hodges (1951) y extendido por Cover y Hart (1967).",
         Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.7),
         size=14, color=BLACK, italic=True, line_spacing=1.3)

add_card(s, Inches(0.3), Inches(2.1), Inches(4.65), Inches(2.2),
         title="Doble proposito",
         body="* Clasificacion: asigna una categoria (p.ej. fraude / no fraude).\n"
              "* Regresion: predice un valor continuo (p.ej. precio de casa).",
         title_size=15, body_size=12, line_spacing=1.35)

add_card(s, Inches(5.05), Inches(2.1), Inches(4.65), Inches(2.2),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Lazy learner",
         body="No construye un modelo explicito durante el entrenamiento. "
              "Memoriza los datos y solo calcula al momento de predecir.",
         title_size=15, body_size=12, line_spacing=1.35)

add_card(s, Inches(0.3), Inches(4.45), Inches(9.4), Inches(1.6),
         fill=SOFT_BLUE,
         title="No parametrico",
         body="No asume ninguna forma de distribucion de los datos. "
              "Esto le permite capturar fronteras de decision muy flexibles, "
              "sin imponer supuestos como linealidad o normalidad.",
         title_size=15, body_size=13, line_spacing=1.4)

add_card(s, Inches(0.3), Inches(6.2), Inches(9.4), Inches(0.85),
         fill=HIGHLIGHT_GREEN,
         body="Idea clave: K-NN predice mirando a los vecinos mas parecidos "
              "del nuevo punto en el conjunto de entrenamiento.",
         body_size=13)
add_footer(s, 2)


# --- Slide 3: La intuicion ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "La intuicion")

add_rect(s, Inches(0.3), Inches(1.5), Inches(9.4), Inches(1.8), NAVY)
add_text(s, "\"Dime con quien andas y te dire quien eres.\"",
         Inches(0.3), Inches(1.5), Inches(9.4), Inches(1.8),
         size=32, color=WHITE, bold=True, italic=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_card(s, Inches(0.3), Inches(3.6), Inches(4.65), Inches(2.5),
         title="Premisa",
         body="Si un nuevo punto se parece a sus vecinos en el espacio de "
              "caracteristicas, probablemente pertenezca a su misma clase.",
         title_size=16, body_size=13, line_spacing=1.4)

add_card(s, Inches(5.05), Inches(3.6), Inches(4.65), Inches(2.5),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Idea central",
         body="Proximidad implica similitud.\n\n"
              "Toda la mecanica del algoritmo se reduce a medir distancias "
              "y consultar a los vecinos.",
         title_size=16, body_size=13, line_spacing=1.4)

add_card(s, Inches(0.3), Inches(6.25), Inches(9.4), Inches(0.85),
         fill=HIGHLIGHT_GREEN,
         body="En politicas publicas: clasificar municipios, hogares o "
              "beneficiarios buscando a quienes mas se le parecen.",
         body_size=13)
add_footer(s, 3)


# --- Slide 4: Espacio de caracteristicas ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "El espacio de caracteristicas")

add_formula_box(s, "x = (x1, x2, ..., xn)",
                "Cada observacion es un vector con n atributos",
                Inches(0.3), Inches(1.3), Inches(9.4), Inches(1.2))

add_card(s, Inches(0.3), Inches(2.8), Inches(4.65), Inches(2.4),
         title="Geometria",
         body="Cada observacion es un punto en un espacio de n dimensiones. "
              "Con 2 variables: plano. Con 3: cubo. Con n: hipercubo.",
         title_size=15, body_size=13, line_spacing=1.4)

add_card(s, Inches(5.05), Inches(2.8), Inches(4.65), Inches(2.4),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Cercania = metrica",
         body='"Cercania" no es algo abstracto: se mide con una formula de '
              "distancia que asigna un numero a cada par de puntos.",
         title_size=15, body_size=13, line_spacing=1.4)

add_card(s, Inches(0.3), Inches(5.35), Inches(9.4), Inches(1.7),
         fill=SOFT_BLUE,
         title="Ejemplo (3 variables)",
         body="Un municipio podria describirse con (PIB per capita, indice "
              "de pobreza, escolaridad promedio). Cada uno es un punto en R3 "
              "y municipios parecidos quedan cerca entre si.",
         title_size=14, body_size=13, line_spacing=1.4)
add_footer(s, 4)


# --- Slide 5: Distancia Euclidiana ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Distancia euclidiana (la mas usada)")

add_formula_box(s,
                "d(p,q) = sqrt( sum_{i=1..n} (pi - qi)^2 )",
                "Generalizacion de Pitagoras a n dimensiones",
                Inches(0.3), Inches(1.25), Inches(9.4), Inches(1.2))

# Imagen
s.shapes.add_picture("img_knn/euclidiana.png",
                     Inches(0.3), Inches(2.7),
                     width=Inches(5.5))

add_card(s, Inches(6.0), Inches(2.7), Inches(3.7), Inches(2.3),
         fill=LIGHT_BLUE,
         title="Notacion",
         body="* p y q: dos puntos en R^n\n"
              "* pi, qi: i-esima coordenada\n"
              "* n: numero de variables",
         title_size=14, body_size=12, line_spacing=1.4)

add_card(s, Inches(6.0), Inches(5.1), Inches(3.7), Inches(1.9),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Caso 2D (Pitagoras)",
         body="d = sqrt( (x2-x1)^2 + (y2-y1)^2 )",
         title_size=13, body_size=12, line_spacing=1.4)
add_footer(s, 5)


# --- Slide 6: Otras metricas (tabla) ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Otras metricas de distancia")

# Tabla manual
col_w = [1.8, 3.5, 4.1]
col_x = [0.3, 2.1, 5.6]
row_h = 0.65
top = 1.4

# Header de tabla
for i, (x, w) in enumerate(zip(col_x, col_w)):
    add_rect(s, Inches(x), Inches(top), Inches(w), Inches(row_h), NAVY)
headers = ["Metrica", "Formula", "Cuando usarla"]
for i, (x, w, h) in enumerate(zip(col_x, col_w, headers)):
    add_text(s, h, Inches(x), Inches(top), Inches(w), Inches(row_h),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("Manhattan (L1)", "sum |pi - qi|",
     "Variables tipo grilla, robusta a outliers"),
    ("Minkowski", "( sum |pi - qi|^p )^(1/p)",
     "Generaliza L1 y L2 (parametro p)"),
    ("Coseno", "1 - (p . q) / (||p|| ||q||)",
     "Texto, datos de alta dimensionalidad"),
    ("Hamming", "# coordenadas distintas",
     "Variables categoricas o binarias"),
]
for r, (m, f, u) in enumerate(rows):
    y = top + row_h * (r + 1)
    fill = SOFT_BLUE if r % 2 == 0 else LIGHT_BLUE
    for i, (x, w) in enumerate(zip(col_x, col_w)):
        add_rect(s, Inches(x), Inches(y), Inches(w), Inches(row_h), fill)
    add_text(s, m, Inches(col_x[0]) + Inches(0.1), Inches(y),
             Inches(col_w[0]) - Inches(0.1), Inches(row_h),
             size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, f, Inches(col_x[1]) + Inches(0.1), Inches(y),
             Inches(col_w[1]) - Inches(0.1), Inches(row_h),
             size=12, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, u, Inches(col_x[2]) + Inches(0.1), Inches(y),
             Inches(col_w[2]) - Inches(0.1), Inches(row_h),
             size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)

add_card(s, Inches(0.3), Inches(5.55), Inches(9.4), Inches(1.5),
         fill=HIGHLIGHT_GREEN,
         title="Como elegir la metrica?",
         body="No hay una respuesta unica: depende del tipo de variables y "
              "del dominio. La euclidiana es el punto de partida por defecto; "
              "Manhattan ayuda con outliers; coseno funciona bien con texto.",
         title_size=13, body_size=12, line_spacing=1.4)
add_footer(s, 6)


# --- Slide 7: Algoritmo paso a paso ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "El algoritmo, paso a paso")

# Paso 0
add_card(s, Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.75),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Paso 0: definir k (entero positivo, idealmente impar en "
               "clasificacion binaria).",
         title_size=14)

# 4 cards horizontales con pasos a-d
steps = [
    ("a) Calcular",
     "Distancia entre el nuevo punto x y cada observacion del conjunto "
     "de entrenamiento."),
    ("b) Ordenar",
     "Ordenar las distancias de menor a mayor."),
    ("c) Seleccionar",
     "Tomar los k vecinos mas cercanos a x."),
    ("d) Predecir",
     "Clasificacion: clase mayoritaria entre los k vecinos.\n"
     "Regresion: promedio del valor objetivo."),
]
card_w = 2.2
gap = 0.1
start_x = 0.3
for i, (t, b) in enumerate(steps):
    x = start_x + i * (card_w + gap)
    fill = LIGHT_BLUE
    add_card(s, Inches(x), Inches(2.1), Inches(card_w), Inches(3.5),
             fill=fill,
             title=t, body=b,
             title_size=14, body_size=11, line_spacing=1.35)
    # Numero circular
    add_rect(s, Inches(x + card_w / 2 - 0.2), Inches(2.0),
             Inches(0.4), Inches(0.4), NAVY)
    add_text(s, str(i + 1), Inches(x + card_w / 2 - 0.2), Inches(2.0),
             Inches(0.4), Inches(0.4), size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_card(s, Inches(0.3), Inches(5.85), Inches(9.4), Inches(1.2),
         fill=HIGHLIGHT_GREEN,
         body="Todo el costo computacional vive en la prediccion: por eso "
              "K-NN es rapido de entrenar pero lento en prediccion sobre "
              "datasets grandes.",
         body_size=13, line_spacing=1.4)
add_footer(s, 7)


# --- Slide 8: Voto mayoritario ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Voto mayoritario (formalizacion)")

add_text(s, "Para clasificacion, la prediccion para un punto x es:",
         Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.4),
         size=14, color=BLACK, italic=True)

add_formula_box(s,
                "y(x) = argmax_{c en C}  sum_{i en Nk(x)}  1(yi = c)",
                "Se asigna la clase mas frecuente entre los k vecinos",
                Inches(0.3), Inches(1.7), Inches(9.4), Inches(1.2))

add_card(s, Inches(0.3), Inches(3.2), Inches(4.65), Inches(2.4),
         title="Notacion",
         body="* Nk(x): conjunto de los k vecinos mas cercanos de x.\n"
              "* C: conjunto de clases posibles.\n"
              "* 1(.): funcion indicadora (vale 1 si se cumple).",
         title_size=15, body_size=12, line_spacing=1.4)

add_card(s, Inches(5.05), Inches(3.2), Inches(4.65), Inches(2.4),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Version ponderada",
         body="Cada voto pesa wi = 1 / di,\n"
              "donde di es la distancia al vecino i.\n\n"
              "Mas peso a los vecinos mas cercanos.",
         title_size=15, body_size=12, line_spacing=1.4)

add_card(s, Inches(0.3), Inches(5.8), Inches(9.4), Inches(1.3),
         fill=HIGHLIGHT_GREEN,
         title="Ventaja del voto ponderado",
         body="Reduce el efecto de vecinos lejanos cuando k es grande y "
              "estabiliza la prediccion cerca de la frontera de decision.",
         title_size=13, body_size=12, line_spacing=1.4)
add_footer(s, 8)


# --- Slide 9: Como elegir k (con imagen) ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Como elegir k?")

s.shapes.add_picture("img_knn/k_efecto.png",
                     Inches(0.3), Inches(1.2),
                     width=Inches(9.4))

add_card(s, Inches(0.3), Inches(5.0), Inches(4.65), Inches(1.2),
         fill=SOFT_RED, title_color=DARK_RED,
         title="k pequeno (p.ej. k = 1)",
         body="Fronteras irregulares, alta varianza, sensible al ruido -> "
              "sobreajuste.",
         title_size=13, body_size=11, line_spacing=1.3)

add_card(s, Inches(5.05), Inches(5.0), Inches(4.65), Inches(1.2),
         fill=SOFT_BLUE, title_color=NAVY,
         title="k grande",
         body="Fronteras suaves, alto sesgo -> riesgo de subajuste.",
         title_size=13, body_size=11, line_spacing=1.3)

add_card(s, Inches(0.3), Inches(6.3), Inches(9.4), Inches(0.8),
         fill=HIGHLIGHT_GREEN,
         body="Reglas practicas: empezar con k ~ sqrt(n), elegir k impar en "
              "binaria, y elegir el k final con validacion cruzada.",
         body_size=12)
add_footer(s, 9)


# --- Slide 10: K-NN visual con k=5 ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Visualizacion de K-NN")

s.shapes.add_picture("img_knn/knn_k5.png",
                     Inches(0.3), Inches(1.2),
                     width=Inches(5.8))

add_card(s, Inches(6.3), Inches(1.2), Inches(3.4), Inches(2.0),
         title="Lectura",
         body="El punto verde (nuevo) consulta a sus 5 vecinos mas cercanos. "
              "El circulo punteado marca el radio.",
         title_size=14, body_size=12, line_spacing=1.4)

add_card(s, Inches(6.3), Inches(3.3), Inches(3.4), Inches(2.0),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Decision",
         body="Si la mayoria de esos 5 vecinos pertenece a la Clase B, "
              "el punto se clasifica como Clase B.",
         title_size=14, body_size=12, line_spacing=1.4)

add_card(s, Inches(6.3), Inches(5.4), Inches(3.4), Inches(1.7),
         fill=SOFT_BLUE,
         title="Importante",
         body="Cambiar k cambia el radio del circulo y, con el, la "
              "prediccion.",
         title_size=13, body_size=12, line_spacing=1.4)
add_footer(s, 10)


# --- Slide 11: Escalar variables ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, Inches(0.3), Inches(7.5), Inches(0.7), DARK_RED)
add_text(s, "Escalar las variables es OBLIGATORIO",
         Inches(0.0), Inches(0.3), Inches(7.5),
         Inches(0.7), size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s,
         "Las variables con escalas grandes dominan el calculo de "
         "distancia.",
         Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.5),
         size=14, color=BLACK, italic=True)

s.shapes.add_picture("img_knn/escalado.png",
                     Inches(0.3), Inches(1.7),
                     width=Inches(9.4))

add_card(s, Inches(0.3), Inches(5.4), Inches(4.65), Inches(1.6),
         fill=LIGHT_BLUE,
         title="Estandarizacion (z-score)",
         body="x' = (x - mu) / sigma\n"
              "Media 0 y desviacion 1.",
         title_size=13, body_size=12, line_spacing=1.4)

add_card(s, Inches(5.05), Inches(5.4), Inches(4.65), Inches(1.6),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Normalizacion min-max",
         body="x' = (x - x_min) / (x_max - x_min)\n"
              "Todo queda en [0, 1].",
         title_size=13, body_size=12, line_spacing=1.4)
add_footer(s, 11)


# --- Slide 12: Ventajas ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Ventajas de K-NN")

ventajas = [
    ("Conceptualmente simple",
     "Faciles de entender e interpretar: 'mira a tus vecinos y decide'. "
     "Util para comunicar a audiencias no tecnicas."),
    ("Sin fase de entrenamiento",
     "No hay funcion de costo que optimizar ni parametros que ajustar; "
     "basta con almacenar los datos."),
    ("Multi-clase natural",
     "Maneja sin esfuerzo problemas con tres o mas categorias: el voto "
     "mayoritario se generaliza directamente."),
    ("No parametrico",
     "No asume ninguna distribucion subyacente; se adapta a la "
     "estructura real de los datos."),
    ("Fronteras no lineales",
     "Captura fronteras de decision arbitrariamente complejas sin tener "
     "que especificarlas a priori."),
    ("Versatil",
     "Sirve tanto para clasificacion como para regresion, con la misma "
     "logica subyacente."),
]
posiciones = [
    (0.3, 1.2), (3.43, 1.2), (6.57, 1.2),
    (0.3, 4.2), (3.43, 4.2), (6.57, 4.2),
]
for (x, y), (t, b) in zip(posiciones, ventajas):
    add_card(s, Inches(x), Inches(y), Inches(3.0), Inches(2.9),
             fill=SOFT_GREEN, title_color=DARK_GREEN,
             title=t, body=b,
             title_size=13, body_size=11, line_spacing=1.35)
add_footer(s, 12)


# --- Slide 13: Desventajas ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, Inches(0.3), Inches(7.5), Inches(0.7), DARK_RED)
add_text(s, "Desventajas de K-NN",
         Inches(0.0), Inches(0.3), Inches(7.5),
         Inches(0.7), size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

desventajas = [
    ("Costoso en prediccion",
     "Complejidad O(n x d) por punto a predecir: se vuelve lento con "
     "muchos datos o muchas variables."),
    ("Maldicion de la dimensionalidad",
     "En alta dimension todas las distancias se vuelven similares y el "
     "concepto de 'vecino cercano' pierde sentido."),
    ("Requiere mucha memoria",
     "Hay que guardar todo el conjunto de entrenamiento; no hay un "
     "modelo compacto que reemplace los datos."),
    ("Sensible a escalas",
     "Una variable mal escalada puede dominar el calculo y arruinar "
     "las predicciones."),
    ("Variables irrelevantes danan",
     "K-NN no distingue por si solo variables informativas de ruido: "
     "todas suman a la distancia."),
    ("Sensible al desbalance",
     "Clases con muchos ejemplos tienden a 'ganar' los votos, sesgando "
     "la prediccion contra las minoritarias."),
]
for (x, y), (t, b) in zip(posiciones, desventajas):
    add_card(s, Inches(x), Inches(y), Inches(3.0), Inches(2.9),
             fill=SOFT_RED, title_color=DARK_RED,
             title=t, body=b,
             title_size=13, body_size=11, line_spacing=1.35)
add_footer(s, 13)


# --- Slide 14: Variantes y mejoras ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Variantes y mejoras")

add_card(s, Inches(0.3), Inches(1.2), Inches(4.65), Inches(2.6),
         fill=LIGHT_BLUE,
         title="Weighted K-NN",
         body="Cada voto se pondera por el inverso de la distancia "
              "(wi = 1/di). Los vecinos mas cercanos influyen mas que "
              "los lejanos, reduciendo el impacto del ruido.",
         title_size=15, body_size=12, line_spacing=1.4)

add_card(s, Inches(5.05), Inches(1.2), Inches(4.65), Inches(2.6),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="KD-Trees / Ball-Trees",
         body="Estructuras de datos que permiten buscar a los k vecinos "
              "en tiempo sub-lineal en lugar de comparar contra todos "
              "los puntos. Muy efectivas en dimensionalidad moderada.",
         title_size=15, body_size=12, line_spacing=1.4)

add_card(s, Inches(0.3), Inches(3.95), Inches(4.65), Inches(2.6),
         fill=SOFT_BLUE,
         title="Approximate Nearest Neighbors (ANN)",
         body="Algoritmos aproximados que sacrifican exactitud por "
              "velocidad: FAISS, Annoy, HNSW. Permiten K-NN sobre "
              "millones o miles de millones de puntos.",
         title_size=15, body_size=12, line_spacing=1.4)

add_card(s, Inches(5.05), Inches(3.95), Inches(4.65), Inches(2.6),
         fill=LIGHT_BLUE,
         title="Reduccion de variables previa",
         body="Aplicar PCA, seleccion por feature importance u otras "
              "tecnicas para reducir la dimensionalidad ANTES de K-NN, "
              "mitigando la maldicion de la dimensionalidad.",
         title_size=15, body_size=12, line_spacing=1.4)

add_card(s, Inches(0.3), Inches(6.7), Inches(9.4), Inches(0.55),
         fill=HIGHLIGHT_GREEN,
         body="Casi todas las mejoras atacan dos problemas: velocidad "
              "de busqueda o calidad de la metrica en alta dimension.",
         body_size=11)
add_footer(s, 14)


# --- Slide 15: Aplicaciones ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Aplicaciones")

aplicaciones = [
    ("Sistemas de recomendacion",
     "Usuarios o productos parecidos -> recomendaciones (Netflix, "
     "Amazon, Spotify)."),
    ("Reconocimiento de patrones",
     "Imagenes, escritura a mano, huellas: clasificar comparando con "
     "ejemplos previos."),
    ("Clasificacion de texto",
     "Documentos cercanos en espacios TF-IDF o embeddings; deteccion "
     "de tema o sentimiento."),
    ("Imputacion de faltantes",
     "Reemplazar valores ausentes por el promedio o moda de los k "
     "vecinos mas parecidos."),
    ("Deteccion de anomalias",
     "Puntos cuya distancia al k-esimo vecino es inusualmente grande "
     "se marcan como atipicos."),
    ("Politica publica",
     "Clasificacion de municipios por perfil socioeconomico, "
     "segmentacion de hogares, unidades territoriales similares."),
]
for (x, y), (t, b) in zip(posiciones, aplicaciones):
    add_card(s, Inches(x), Inches(y), Inches(3.0), Inches(2.9),
             fill=LIGHT_BLUE,
             title=t, body=b,
             title_size=13, body_size=11, line_spacing=1.35)
add_footer(s, 15)


# --- Slide 16: Resumen ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Resumen")

resumen_puntos = [
    ("1. Clasificacion por cercania",
     "K-NN clasifica un punto usando los k puntos mas cercanos del "
     "conjunto de entrenamiento."),
    ("2. Distancia euclidiana por defecto",
     "Es la metrica mas comun, pero no la unica: Manhattan, Minkowski, "
     "coseno o Hamming segun el problema."),
    ("3. k es un hiperparametro critico",
     "k pequeno sobreajusta; k grande subajusta. Se elige con "
     "validacion cruzada."),
    ("4. Siempre escalar las variables",
     "Estandarizacion o normalizacion antes de aplicar K-NN: las "
     "escalas grandes dominan el calculo de distancia."),
]
y0 = 1.2
for i, (t, b) in enumerate(resumen_puntos):
    add_card(s, Inches(0.3), Inches(y0 + i * 1.15), Inches(9.4),
             Inches(1.05),
             fill=LIGHT_BLUE,
             title=t, body=b,
             title_size=14, body_size=12, line_spacing=1.3)

add_card(s, Inches(0.3), Inches(5.95), Inches(9.4), Inches(1.1),
         fill=HIGHLIGHT_GREEN,
         title="Cuando usarlo",
         body="Ideal cuando se necesita una frontera de decision flexible "
              "y se dispone de relativamente pocos datos bien preprocesados.",
         title_size=14, body_size=12, line_spacing=1.4)
add_footer(s, 16)


# --- Slide 17: Para reflexionar ---
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Para reflexionar")

add_text(s, "Tres preguntas para llevarse a casa:",
         Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.5),
         size=15, color=BLACK, italic=True)

preguntas = [
    ("?",
     "Que metrica de distancia tendria sentido para tu dataset?",
     "Piensa en el tipo de variables (continuas, binarias, "
     "categoricas, texto) y en si hay outliers."),
    ("?",
     "Como afectaria tener clases desbalanceadas a las predicciones?",
     "Que vecinos terminarian dominando el voto y como podrias "
     "compensar esa asimetria?"),
    ("?",
     "En que casos K-NN seria preferible a una regresion logistica "
     "o un arbol de decision?",
     "Pista: tamano del dataset, forma de las fronteras, necesidad "
     "de interpretabilidad y velocidad de prediccion."),
]
y0 = 1.9
for i, (num, p, hint) in enumerate(preguntas):
    yi = y0 + i * 1.6
    # Badge circular
    add_rect(s, Inches(0.3), Inches(yi + 0.15),
             Inches(0.55), Inches(0.55), NAVY)
    add_text(s, str(i + 1), Inches(0.3), Inches(yi + 0.15),
             Inches(0.55), Inches(0.55), size=22, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Card
    add_card(s, Inches(1.05), Inches(yi), Inches(8.65), Inches(1.4),
             fill=SOFT_BLUE,
             title=p, body=hint,
             title_size=14, body_size=11, line_spacing=1.35)
add_footer(s, 17)


# --- Slide Cierre ---
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_text(s, "Gracias", Inches(0.6), Inches(1.7),
         Inches(8.8), Inches(2), size=96, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_rect(s, Inches(3.0), Inches(5.4), Inches(4.0), Inches(1.0), WHITE)
add_text(s, "Preguntas?", Inches(3.0), Inches(5.4),
         Inches(4.0), Inches(1.0), size=28, color=NAVY, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


OUT = "K-NN_Vecinos_Cercanos.pptx"
prs.save(OUT)
print(f"PPTX guardado: {OUT}")
print(f"Slides totales: {len(prs.slides)}")
